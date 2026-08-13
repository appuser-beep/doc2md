# -*- coding: utf-8 -*-
"""跨格式转换保真回归：正文/表头/特殊字符不得丢失或拆坏。"""

from __future__ import annotations

import sys
import tempfile
import unittest
from datetime import date, datetime
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

from cleanup import clean_markdown_light  # noqa: E402
from converter import convert_path  # noqa: E402
from excel_convert import convert_excel_to_markdown  # noqa: E402


class TestCleanupTableRepair(unittest.TestCase):
    def test_word_empty_header_promoted(self):
        raw = "|  |  |\n| --- | --- |\n| 姓名 | 部门 |\n| 李四 | 研发 |\n"
        md = clean_markdown_light(raw)
        self.assertIn("| 姓名 | 部门 |", md)
        self.assertNotRegex(md.splitlines()[0], r"^\|\s*\|\s*\|$")

    def test_pipe_inside_cell_escaped(self):
        raw = "| 姓名 | 部门 |\n| --- | --- |\n| 李四 | 研发|一组 |\n"
        md = clean_markdown_light(raw)
        self.assertIn("研发\\|一组", md)
        data = [ln for ln in md.splitlines() if ln.startswith("|") and "---" not in ln][-1]
        self.assertIn("\\|", data)

    def test_code_fence_tables_untouched(self):
        raw = (
            "Intro\n\n"
            "```text\n"
            "| keep | me |\n"
            "| --- | --- |\n"
            "| a | b |\n"
            "```\n\n"
            "|  |  |\n"
            "| --- | --- |\n"
            "| 真表头 | 列B |\n"
            "| 1 | 2 |\n"
        )
        md = clean_markdown_light(raw)
        self.assertIn("| keep | me |", md)
        self.assertIn("| a | b |", md)
        self.assertIn("| 真表头 | 列B |", md)
        # 代码块内不应被转义改写
        fence_body = md.split("```text", 1)[1].split("```", 1)[0]
        self.assertNotIn("\\|", fence_body)

    def test_escape_cell_idempotent(self):
        from cleanup import _escape_cell

        once = _escape_cell("研发|一组")
        twice = _escape_cell(once)
        self.assertEqual(once, "研发\\|一组")
        self.assertEqual(twice, once)

    def test_fit_mid_pipe_prefers_short_fragments(self):
        from cleanup import _fit_row_width

        # 中间被 | 劈开的短姓名：合并前两段
        self.assertEqual(_fit_row_width(["张", "三", "研发"], 2), ["张|三", "研发"])
        # 末列含 |：同分时靠后合并，保持首列完整
        self.assertEqual(_fit_row_width(["李四", "研发", "一组"], 2), ["李四", "研发|一组"])
        self.assertEqual(_fit_row_width(["A", "B", "C", "D"], 3), ["A", "B", "C|D"])

    def test_mid_pipe_table_cleanup(self):
        from cleanup import _split_table_row

        raw = "| 姓名 | 部门 |\n| --- | --- |\n| 张|三 | 研发 |\n"
        md = clean_markdown_light(raw)
        self.assertIn("张\\|三", md)
        self.assertIn("研发", md)
        data = [ln for ln in md.splitlines() if ln.startswith("|") and "---" not in ln][-1]
        cells = _split_table_row(data)
        self.assertEqual(len(cells), 2)
        self.assertEqual(cells[0].replace("\\", ""), "张|三")

    def test_short_sep_does_not_collapse_columns(self):
        raw = (
            "| 列A | 列B | 列C |\n"
            "| --- | --- |\n"
            "| 1 | 2 | 3 |\n"
            "| 4 | 5 | 6 |\n"
        )
        md = clean_markdown_light(raw)
        data = [ln for ln in md.splitlines() if ln.startswith("|") and "---" not in ln]
        self.assertGreaterEqual(len(data), 2)
        from cleanup import _split_table_row

        self.assertEqual(len(_split_table_row(data[0])), 3)
        self.assertEqual(len(_split_table_row(data[-1])), 3)
        self.assertIn("列C", md)

    def test_aligned_table_preserved_in_light(self):
        raw = (
            "| 左 | 中 | 右 |\n"
            "| :--- | :---: | ---: |\n"
            "| a | b | c |\n"
        )
        md = clean_markdown_light(raw)
        self.assertIn(":---", md)
        self.assertIn(":---:", md)
        self.assertIn("---:", md)

    def test_short_sep_full_clean_no_junk_row(self):
        from cleanup import clean_markdown

        raw = (
            "| A | B | C |\n"
            "| --- | --- |\n"
            "| 1 | 2 | 3 |\n"
        )
        md = clean_markdown(raw)
        lines = [ln for ln in md.splitlines() if ln.startswith("|")]
        # 不应出现把短分隔补空后当成数据行
        self.assertFalse(any(ln.count("|") >= 4 and "---" in ln and ln.strip().endswith("|  |") for ln in lines))
        from cleanup import _split_table_row

        data = [ln for ln in lines if "---" not in ln]
        self.assertEqual(len(_split_table_row(data[0])), 3)
        self.assertEqual(len(_split_table_row(data[-1])), 3)

    def test_image_alt_with_bracket_repaired(self):
        from cleanup import clean_markdown_light

        raw = "![C:\\tmp\\foo]$bar.png](data:image/png;base64,aaa)\n"
        md = clean_markdown_light(raw)
        self.assertIn("](data:image/png;base64,aaa)", md)
        # alt 内不应再残留未转义的 ]
        alt = md.split("![", 1)[1].split("](", 1)[0]
        self.assertNotIn("]", alt)


class TestExcelFidelity(unittest.TestCase):
    def test_date_percent_currency_and_pipe(self):
        from openpyxl import Workbook

        with tempfile.TemporaryDirectory() as td:
            path = Path(td) / "types.xlsx"
            wb = Workbook()
            ws = wb.active
            ws["A1"] = "日期"
            ws["B1"] = date(2024, 5, 1)
            ws["A2"] = "时间"
            ws["B2"] = datetime(2024, 5, 1, 15, 30)
            ws["A3"] = "比例"
            ws["B3"] = 0.15
            ws["B3"].number_format = "0%"
            ws["A4"] = "部门"
            ws["B4"] = "研发|一组"
            ws["A5"] = "金额"
            ws["B5"] = 1234.5
            ws["B5"].number_format = '¥#,##0.00'
            ws["A6"] = "区位货币"
            ws["B6"] = 12.5
            ws["B6"].number_format = "[$¥-804]#,##0.00"
            ws["A7"] = "科学"
            ws["B7"] = 12345.0
            ws["B7"].number_format = "0.00E+00"
            ws["A8"] = "千分位"
            ws["B8"] = 12345.6
            ws["B8"].number_format = "#,##0.00"
            wb.save(path)
            md = convert_excel_to_markdown(path)
            self.assertIn("2024-05-01", md)
            self.assertNotIn("2024-05-01 00:00:00", md)
            self.assertIn("2024-05-01 15:30:00", md)
            self.assertIn("15%", md)
            self.assertIn("研发\\|一组", md)
            self.assertIn("¥", md)
            self.assertIn("1,234.50", md)
            self.assertIn("12.50", md)
            self.assertRegex(md, r"1\.234500E\+04")
            self.assertIn("12,345.60", md)
            # 千分位小数位按格式，而非误把 .000 当成 .00
            from excel_convert import _format_number

            self.assertEqual(_format_number(1.5, "#,##0.0"), "1.5")
            self.assertEqual(_format_number(1.23456, "#,##0.000"), "1.235")


class TestPostcheck(unittest.TestCase):
    def test_pdf_empty_and_short(self):
        from converter import postcheck_result

        tip = postcheck_result("scan.pdf", "")
        self.assertIsNotNone(tip)
        self.assertIn("扫描件", tip)
        tip2 = postcheck_result("scan.pdf", " ![](x) \n ")
        self.assertIsNotNone(tip2)
        tip3 = postcheck_result("ok.pdf", "PDF_TEXT_MARKER Revenue 1280")
        self.assertIsNone(tip3)


class TestOfficeFidelity(unittest.TestCase):
    def test_docx_table_header_and_body(self):
        from docx import Document

        with tempfile.TemporaryDirectory() as td:
            path = Path(td) / "t.docx"
            d = Document()
            d.add_heading("DOCX_TITLE_MARKER", level=1)
            d.add_paragraph("正文 ALPHA_WORD_TEXT")
            t = d.add_table(rows=2, cols=2)
            t.rows[0].cells[0].text = "姓名"
            t.rows[0].cells[1].text = "部门"
            t.rows[1].cells[0].text = "李四"
            t.rows[1].cells[1].text = "研发|一组"
            d.sections[0].header.paragraphs[0].text = "页眉HEADER_MARK"
            d.sections[0].footer.paragraphs[0].text = "页脚FOOTER_MARK"
            d.save(path)
            md = convert_path(str(path), local_only=False)
            self.assertIn("DOCX_TITLE_MARKER", md)
            self.assertIn("ALPHA_WORD_TEXT", md)
            self.assertIn("姓名", md)
            self.assertIn("李四", md)
            self.assertIn("HEADER_MARK", md)
            self.assertIn("FOOTER_MARK", md)
            # 表头不应再是空行
            lines = [ln for ln in md.splitlines() if ln.startswith("|")]
            self.assertTrue(lines)
            self.assertRegex(lines[0], r"姓名")

    def test_pptx_title_body_table(self):
        from pptx import Presentation
        from pptx.util import Inches

        with tempfile.TemporaryDirectory() as td:
            path = Path(td) / "t.pptx"
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "PPT_TITLE_MARKER"
            slide.placeholders[1].text = "PPT_BODY_TEXT"
            table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(5), Inches(1)).table
            table.cell(0, 0).text = "列A"
            table.cell(0, 1).text = "列B"
            table.cell(1, 0).text = "值1"
            table.cell(1, 1).text = "值|2"
            prs.save(path)
            md = convert_path(str(path), local_only=False)
            self.assertIn("PPT_TITLE_MARKER", md)
            self.assertIn("PPT_BODY_TEXT", md)
            self.assertIn("列A", md)
            self.assertIn("值1", md)

    def test_pdf_text(self):
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4

        with tempfile.TemporaryDirectory() as td:
            path = Path(td) / "t.pdf"
            c = canvas.Canvas(str(path), pagesize=A4)
            c.drawString(72, 720, "PDF_TEXT_MARKER Revenue 1280")
            c.save()
            md = convert_path(str(path), local_only=False)
            self.assertIn("PDF_TEXT_MARKER", md)
            self.assertIn("1280", md)

    def test_html_csv_json(self):
        with tempfile.TemporaryDirectory() as td:
            root = Path(td)
            (root / "a.html").write_text(
                "<html><body><h1>HTML_H1</h1><p>华北 1280</p>"
                "<table><tr><th>区</th><th>值</th></tr>"
                "<tr><td>华北</td><td>1280</td></tr></table></body></html>",
                encoding="utf-8",
            )
            (root / "a.csv").write_text("区,值\n华北,1280\n", encoding="utf-8")
            (root / "a.json").write_text(
                '{"mark":"JSON_MARK","region":"华北","value":1280}',
                encoding="utf-8",
            )
            html = convert_path(str(root / "a.html"), local_only=False)
            csv = convert_path(str(root / "a.csv"), local_only=False)
            js = convert_path(str(root / "a.json"), local_only=False)
            self.assertIn("HTML_H1", html)
            self.assertIn("1280", html)
            self.assertIn("华北", csv)
            self.assertIn("JSON_MARK", js)


class TestLlmUsabilityFixes(unittest.TestCase):
    def test_real_zip_cn_names_if_present(self):
        zpath = Path(r"C:\Users\wluser\Desktop\档案整理活动.zip")
        if not zpath.is_file():
            self.skipTest("桌面样例 ZIP 不存在")
        md = convert_path(str(zpath), local_only=False)
        self.assertIn("档案", md)
        # 文件名不应再是 CP437 乱码
        self.assertNotRegex(md, r"## 文件：`[^`]*[╡░╕╒√└φ╗][^`]*`")

    def test_real_excel_pair_cols_if_present(self):
        xpath = Path(
            r"C:\Users\wluser\Documents\WPSDrive\1359559257\WPS云盘\国奖分享会在线表格.xlsx"
        )
        if not xpath.is_file():
            self.skipTest("国奖表格不存在")
        md = convert_path(str(xpath), local_only=False)
        self.assertIn("姓名", md)
        self.assertNotIn("| 年级专业 | 年级专业 |", md)
        self.assertNotIn("| 学号 | 学号 |", md)


class TestIpynbZipExtras(unittest.TestCase):
    def test_ipynb_markdown_output(self):
        import json

        from ipynb_convert import convert_ipynb_to_markdown

        nb = {
            "cells": [
                {
                    "cell_type": "code",
                    "source": ["1"],
                    "outputs": [
                        {
                            "output_type": "execute_result",
                            "data": {
                                "text/markdown": ["## MD_OUT_MARK\n", "body"],
                                "text/plain": ["plain should lose"],
                            },
                            "metadata": {},
                            "execution_count": 1,
                        }
                    ],
                }
            ],
            "metadata": {},
            "nbformat": 4,
            "nbformat_minor": 5,
        }
        with tempfile.TemporaryDirectory() as td:
            path = Path(td) / "t.ipynb"
            path.write_text(json.dumps(nb), encoding="utf-8")
            md = convert_ipynb_to_markdown(path)
            self.assertIn("MD_OUT_MARK", md)
            self.assertNotIn("plain should lose", md)

    def test_ipynb_markdown_with_image(self):
        import json

        from ipynb_convert import convert_ipynb_to_markdown

        nb = {
            "cells": [
                {
                    "cell_type": "code",
                    "source": ["plot()"],
                    "outputs": [
                        {
                            "output_type": "display_data",
                            "data": {
                                "text/markdown": ["## CHART_MD_MARK"],
                                "image/png": "aaaa",
                            },
                            "metadata": {},
                        }
                    ],
                }
            ],
            "metadata": {},
            "nbformat": 4,
            "nbformat_minor": 5,
        }
        with tempfile.TemporaryDirectory() as td:
            path = Path(td) / "img.ipynb"
            path.write_text(json.dumps(nb), encoding="utf-8")
            md = convert_ipynb_to_markdown(path)
            self.assertIn("CHART_MD_MARK", md)
            self.assertIn("embedded-image", md)

    def test_zip_skips_min_js(self):
        from zip_convert import _should_skip_name

        self.assertTrue(_should_skip_name("assets/app.min.js"))
        self.assertTrue(_should_skip_name("bundle.js.map"))
        self.assertFalse(_should_skip_name("src/app.js"))

    def test_zip_skips_path_traversal_names(self):
        from zip_convert import _is_unsafe_zip_member, _should_skip_name

        self.assertTrue(_is_unsafe_zip_member("../evil.txt"))
        self.assertTrue(_is_unsafe_zip_member("a/../../b.txt"))
        self.assertTrue(_is_unsafe_zip_member("C:/Windows/x.txt"))
        self.assertTrue(_should_skip_name("../evil.txt"))
        self.assertFalse(_is_unsafe_zip_member("docs/说明.txt"))


if __name__ == "__main__":
    unittest.main(verbosity=2)
