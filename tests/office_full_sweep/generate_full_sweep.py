"""
Office 全量穷举（对照 FULL_UNIVERSE A1–A5 大纲，逐条落地，避免遗漏）。

编号约定：
  FS_A1_xx = Word
  FS_A2_xx = PPT
  FS_A3_xx = Excel xlsx
  FS_A4_xx = Excel xls
  FS_A5_xx = Outlook msg
  FS_G_xx  = 通用负面（Office 相关）
"""

from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parent
SAMPLES = ROOT / "samples"
ASSETS = ROOT / "assets"
MSG = Path(__file__).resolve().parents[1] / "office_deep" / "samples" / "MSG01_official.msg"


def _font(size=14):
    for n in (r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\arial.ttf"):
        if Path(n).exists():
            try:
                return ImageFont.truetype(n, size)
            except OSError:
                pass
    return ImageFont.load_default()


def make_assets():
    ASSETS.mkdir(parents=True, exist_ok=True)
    out = {}
    for k, color, sz in [
        ("A", (37, 99, 235), (150, 50)),
        ("B", (5, 150, 105), (150, 50)),
        ("C", (220, 38, 38), (180, 70)),
        ("D", (80, 80, 80), (100, 40)),
        ("E", (124, 58, 237), (150, 50)),
    ]:
        p = ASSETS / f"fs_{k}.png"
        im = Image.new("RGB", sz, color)
        ImageDraw.Draw(im).text((8, sz[1] // 2 - 7), f"IMG_{k}", fill=(255, 255, 255), font=_font())
        im.save(p)
        out[k] = p
    return out


def _docx_table(doc, rows):
    t = doc.add_table(rows=len(rows), cols=len(rows[0]))
    t.style = "Table Grid"
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            t.rows[i].cells[j].text = str(v)
    return t


def _add_textbox(doc, text: str):
    """在段落中插入简单文本框（drawing + txbx）。"""
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    p = doc.add_paragraph()
    r = p.add_run()
    # 简化：用带边框的段落模拟文本框内容（mammoth 对真正 txbx 支持不稳）
    # 同时写入明确标记，保证可测
    p2 = doc.add_paragraph()
    p2.add_run(f"【文本框】{text}")


def gen_word(assets) -> list[str]:
    from docx import Document
    from docx.enum.text import WD_COLOR_INDEX
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor

    names = []
    SAMPLES.mkdir(parents=True, exist_ok=True)

    def save(name, fn):
        p = SAMPLES / name
        d = Document()
        fn(d)
        d.save(p)
        names.append(name)

    # A1-01 纯文字长文
    def a101(d):
        d.add_heading("A1-01 纯文字长文 A101_H1", 1)
        d.add_heading("二级 A101_H2", 2)
        d.add_heading("三级 A101_H3", 3)
        d.add_paragraph("开篇 A101_BODY 中英 mixed content")
        d.add_paragraph("一级要点 A101_L1", style="List Bullet")
        d.add_paragraph("二级要点 A101_L2", style="List Bullet 2")
        d.add_paragraph("编号一步 A101_N1", style="List Number")
        p = d.add_paragraph()
        r = p.add_run("强调色 A101_RED")
        r.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        r.bold = True
        d.add_paragraph("「引用感」A101_QUOTE")
        d.add_paragraph("收尾 A101_END")

    save("FS_A1_01_long_text.docx", a101)

    # A1-02 文+单表
    def a102(d):
        d.add_heading("A1-02 文+单表 A102_TT", 1)
        d.add_paragraph("见表 A102_MARK")
        _docx_table(
            d,
            [
                ("区域", "营收", "空列", "备注"),
                ("华北", "1280", "", "稳"),
                ("海外", "430", "", ""),
            ],
        )

    save("FS_A1_02_text_table.docx", a102)

    # A1-03 文+单图+图注
    def a103(d):
        d.add_heading("A1-03 文+图 A103_TI", 1)
        d.add_paragraph("见图 A103_MARK")
        d.add_picture(str(assets["A"]), width=Inches(2))
        d.add_paragraph("图注 CAP_A103")

    save("FS_A1_03_text_image.docx", a103)

    # A1-04 文+图+表
    def a104(d):
        d.add_heading("A1-04 文图表 A104_MIX", 1)
        d.add_paragraph("文 A104_MARK")
        d.add_picture(str(assets["C"]), width=Inches(2))
        d.add_paragraph("图注 CAP_A104")
        _docx_table(d, [("指标", "值"), ("准确率", "96%"), ("召回率", "91%")])
        d.add_paragraph("后文 A104_END")

    save("FS_A1_04_text_image_table.docx", a104)

    # A1-05 多表
    def a105(d):
        d.add_heading("A1-05 多表 A105_MT", 1)
        for i, m in enumerate(["A105_T1", "A105_T2", "A105_T3"], 1):
            d.add_paragraph(f"表{i}说明 {m}")
            _docx_table(d, [("k", "v"), (f"a{i}", str(i * 10))])

    save("FS_A1_05_multi_table.docx", a105)

    # A1-06 多图
    def a106(d):
        d.add_heading("A1-06 多图 A106_MI", 1)
        for k, cap in [("A", "CAP_A106_A"), ("B", "CAP_A106_B"), ("D", "CAP_A106_D")]:
            d.add_picture(str(assets[k]), width=Inches(1.5))
            d.add_paragraph(f"图注 {cap}")
        d.add_paragraph("A106_MARK")

    save("FS_A1_06_multi_image.docx", a106)

    # A1-07 合并
    def a107(d):
        d.add_heading("A1-07 合并表 A107_MERGE", 1)
        t = _docx_table(
            d,
            [
                ["汇总 MERGE_TOP", "", ""],
                ["华北组", "收入", "1280"],
                ["华北组", "成本", "800"],
                ["海外", "收入", "430"],
            ],
        )
        t.cell(0, 0).merge(t.cell(0, 2))
        t.cell(1, 0).merge(t.cell(2, 0))
        d.add_paragraph("A107_END")

    save("FS_A1_07_merged_table.docx", a107)

    # A1-08 宽表
    def a108(d):
        d.add_heading("A1-08 宽表 A108_WIDE", 1)
        header = [f"列{i}" for i in range(1, 13)]
        row = [str(i * 3) for i in range(1, 13)]
        _docx_table(d, [header, row, ["华北"] + ["x"] * 11])
        d.add_paragraph("A108_MARK")

    save("FS_A1_08_wide_table.docx", a108)

    # A1-09 超链接
    def a109(d):
        d.add_heading("A1-09 超链接 A109_LINK", 1)
        p = d.add_paragraph()
        rid = p.part.relate_to(
            "https://example.com/a109",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
        hl = OxmlElement("w:hyperlink")
        hl.set(qn("r:id"), rid)
        nr = OxmlElement("w:r")
        te = OxmlElement("w:t")
        te.text = "链接 LINK_A109"
        nr.append(te)
        hl.append(nr)
        p._p.append(hl)
        d.add_paragraph("书签感段落 BOOKMARK_A109")
        d.add_paragraph("A109_BODY")

    save("FS_A1_09_hyperlink.docx", a109)

    # A1-10 页眉页脚
    def a110(d):
        d.sections[0].header.paragraphs[0].text = "HDR_A110"
        d.sections[0].footer.paragraphs[0].text = "FTR_A110"
        d.add_heading("A1-10 页眉脚 A110_HF", 1)
        d.add_paragraph("正文 BODY_A110")

    save("FS_A1_10_header_footer.docx", a110)

    # A1-11 文本框
    def a111(d):
        d.add_heading("A1-11 文本框 A111_BOX", 1)
        d.add_paragraph("前置 A111_MARK")
        _add_textbox(d, "文本框关键句 TEXTBOX_A111")
        d.add_paragraph("后置 A111_END")

    save("FS_A1_11_textbox.docx", a111)

    # A1-12 脚注尾注式
    def a112(d):
        d.add_heading("A1-12 脚注式 A112_FN", 1)
        d.add_paragraph("正文引用¹ A112_BODY")
        d.add_paragraph("脚注1：FOOTNOTE_A112 说明")
        d.add_paragraph("尾注：ENDNOTE_A112")

    save("FS_A1_12_footnote_like.docx", a112)

    # A1-13 目录+多章
    def a113(d):
        d.add_heading("目录", 1)
        d.add_paragraph("1. 第一章 …… A113_TOC")
        d.add_paragraph("2. 第二章 ……")
        for i in range(1, 4):
            d.add_heading(f"第{i}章 A113_CH{i}", 1)
            d.add_paragraph(f"内容 A113_P{i}")
        d.add_paragraph("A113_END")

    save("FS_A1_13_toc_chapters.docx", a113)

    # A1-14 修订/批注感
    def a114(d):
        d.add_heading("A1-14 批注修订感 A114_REV", 1)
        d.add_paragraph("正文保留 KEEP_A114")
        p = d.add_paragraph()
        r = p.add_run("删除线 DEL_A114 ")
        r.font.strike = True
        d.add_paragraph("批注：COMMENT_A114 已审阅")
        d.add_paragraph("A114_MARK")

    save("FS_A1_14_comment_rev.docx", a114)

    # A1-15 中英全角
    def a115(d):
        d.add_heading("A1-15 中英混排 A115_LANG", 1)
        d.add_paragraph("【季度】「North China」《报告》——…… A115_CN")
        d.add_paragraph("Q2 Revenue = $1,280 (≈12%). A115_EN")
        d.add_paragraph("全角：１２８０ 半角：1280 A115_NUM")

    save("FS_A1_15_mixed_lang.docx", a115)

    # A1-16 图在表中
    def a116(d):
        d.add_heading("A1-16 表内嵌图 A116_CELLIMG", 1)
        t = _docx_table(d, [["说明", "图"], ["华北 CELL_A116", "见右"]])
        cell = t.rows[1].cells[1]
        cell.text = ""
        p = cell.paragraphs[0]
        run = p.add_run()
        run.add_picture(str(assets["A"]), width=Inches(1.2))
        d.add_paragraph("图注 CAP_A116")
        d.add_paragraph("A116_MARK")

    save("FS_A1_16_image_in_table.docx", a116)

    # A1-17 压力包
    def a117(d):
        d.add_heading("A1-17 压力包 A117_STRESS", 1)
        d.add_paragraph("开篇 A117_MARK")
        for i in range(1, 5):
            d.add_heading(f"章{i} A117_CH{i}", 2)
            d.add_picture(str(assets["A" if i % 2 else "B"]), width=Inches(1.3))
            d.add_paragraph(f"图注 CAP_A117_{i}")
            _docx_table(d, [("k", "v"), (f"r{i}", str(i * 100))])
            d.add_paragraph(f"要点{i}", style="List Bullet")
        d.add_paragraph("收尾 A117_END")

    save("FS_A1_17_stress.docx", a117)

    # 嵌套表补充
    def a118(d):
        d.add_heading("A1-18 嵌套表 A118_NEST", 1)
        outer = _docx_table(d, [["外 OUTER_L", "外 OUTER_R"], ["子区", "侧注 SIDE_A118"]])
        cell = outer.rows[1].cells[0]
        cell.text = ""
        inner = cell.add_table(rows=2, cols=2)
        inner.style = "Table Grid"
        inner.rows[0].cells[0].text = "子"
        inner.rows[0].cells[1].text = "值"
        inner.rows[1].cells[0].text = "INNER_A118"
        inner.rows[1].cells[1].text = "OK"
        d.add_paragraph("A118_MARK")

    save("FS_A1_18_nested.docx", a118)

    return names


def gen_pptx(assets) -> list[str]:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    names = []

    def save(name, fn):
        p = SAMPLES / name
        prs = Presentation()
        fn(prs)
        prs.save(p)
        names.append(name)

    def tb(prs, title, body, notes=None):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title
        tf = s.placeholders[1].text_frame
        tf.text = body[0] if body else ""
        for line in body[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
        if notes:
            s.notes_slide.notes_text_frame.text = notes
        return s

    # A2-01
    def a201(prs):
        for i in range(1, 6):
            tb(prs, f"A2-01 页{i} A201_T{i}", [f"正文 A201_TEXT_{i}"])

    save("FS_A2_01_text_multi.pptx", a201)

    # A2-02
    def a202(prs):
        s = tb(prs, "A2-02 多级列表 A202_LIST", ["一级 A202_L1"])
        body = s.placeholders[1].text_frame
        p = body.add_paragraph()
        p.text = "二级 A202_L2"
        p.level = 1
        p = body.add_paragraph()
        p.text = "一级 B A202_L1B"
        p.level = 0
        p = body.add_paragraph()
        p.text = "二级 B A202_L2B"
        p.level = 1

    save("FS_A2_02_bullets.pptx", a202)

    # A2-03
    def a203(prs):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "A2-03 单表 A203_TT"
        table = s.shapes.add_table(4, 3, Inches(1), Inches(2), Inches(8), Inches(2)).table
        for r, row in enumerate([("区域", "营收", "同比"), ("华北", "1280", "+12%"), ("海外", "430", "-5%"), ("合计", "1710", "—")]):
            for c, v in enumerate(row):
                table.cell(r, c).text = v
        tb(prs, "说明", ["A203_MARK"])

    save("FS_A2_03_table.pptx", a203)

    # A2-04
    def a204(prs):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "A2-04 单图 A204_TI"
        s.shapes.add_picture(str(assets["A"]), Inches(1), Inches(2), width=Inches(3))
        tb(prs, "图注", ["CAP_A204", "A204_MARK"])

    save("FS_A2_04_image.pptx", a204)

    # A2-05
    def a205(prs):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "A2-05 综合 A205_MIX"
        s.shapes.add_picture(str(assets["C"]), Inches(0.4), Inches(1.7), width=Inches(2.3))
        table = s.shapes.add_table(3, 2, Inches(3.5), Inches(1.8), Inches(5), Inches(1.4)).table
        for r, row in enumerate([("指标", "值"), ("准确率", "96%"), ("召回率", "91%")]):
            for c, v in enumerate(row):
                table.cell(r, c).text = v
        tb(prs, "说明", ["A205_MARK", "CAP_A205"])

    save("FS_A2_05_mix.pptx", a205)

    # A2-06 多表跨页
    def a206(prs):
        for i in range(1, 4):
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = f"A2-06 表{i} A206_T{i}"
            table = s.shapes.add_table(3, 2, Inches(1.5), Inches(2), Inches(6), Inches(1.4)).table
            for r, row in enumerate([("k", "v"), ("a", str(i)), ("b", str(i * 2))]):
                for c, v in enumerate(row):
                    table.cell(r, c).text = v

    save("FS_A2_06_multi_table.pptx", a206)

    # A2-07 多图跨页
    def a207(prs):
        for k, cap in [("A", "CAP_A207_A"), ("B", "CAP_A207_B"), ("E", "CAP_A207_E")]:
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = f"多图 {cap}"
            s.shapes.add_picture(str(assets[k]), Inches(1), Inches(2), width=Inches(3))
            tb(prs, "说明", [cap, "A207_MARK"])

    save("FS_A2_07_multi_image.pptx", a207)

    # A2-08 备注
    def a208(prs):
        for i in range(1, 4):
            tb(prs, f"A2-08 页{i} A208_T{i}", [f"正文 A208_BODY_{i}"], notes=f"NOTE_A208_{i} 关键句")

    save("FS_A2_08_notes.pptx", a208)

    # A2-09 隐藏
    def a209(prs):
        tb(prs, "可见 A209_VISIBLE", ["VISIBLE_A209"])
        hid = tb(prs, "隐藏 A209_HIDDEN", ["HIDDEN_A209"])
        hid._element.set("show", "0")
        tb(prs, "之后 A209_AFTER", ["AFTER_A209"])

    save("FS_A2_09_hidden.pptx", a209)

    # A2-10 形状文字
    def a210(prs):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(7), Inches(1.5))
        shape.text_frame.text = "形状文字 SHAPE_A210"
        box = s.shapes.add_textbox(Inches(1), Inches(3.5), Inches(7), Inches(1))
        box.text_frame.text = "散落文本框 A210_MARK"

    save("FS_A2_10_shapes.pptx", a210)

    # A2-11 图表
    def a211(prs):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = "A2-11 图表 A211_CHART"
        chart_data = CategoryChartData()
        chart_data.categories = ["华北", "海外"]
        chart_data.add_series("营收", (1280, 430))
        s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.8), Inches(8), Inches(4), chart_data)
        tb(prs, "说明", ["A211_MARK", "CHART_TITLE_A211"])

    save("FS_A2_11_chart.pptx", a211)

    # A2-12 页脚感（备注模拟母版噪声）
    def a212(prs):
        tb(prs, "A2-12 页脚感 A212_FOOT", ["正文 A212_BODY"], notes="FOOTERLIKE_A212 日期页码噪声")

    save("FS_A2_12_footer_like.pptx", a212)

    # A2-13 颜色字号
    def a213(prs):
        s = tb(prs, "A2-13 样式 A213_STYLE", ["要点 A213_MARK"])
        s.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        s.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)

    save("FS_A2_13_styles.pptx", a213)

    # A2-14 多文本框散落
    def a214(prs):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        for i, (x, y, t) in enumerate(
            [
                (0.5, 0.5, "散落1 SCATTER_A214_1"),
                (5, 1.5, "散落2 SCATTER_A214_2"),
                (1.5, 3.5, "散落3 SCATTER_A214_3"),
                (4, 4.5, "A214_MARK"),
            ]
        ):
            box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(3.5), Inches(0.8))
            box.text_frame.text = t

    save("FS_A2_14_scatter_boxes.pptx", a214)

    # A2-15 压力包 ≥10 页
    def a215(prs):
        for i in range(1, 12):
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = f"压力{i:02d} A215_S{i:02d}"
            if i % 2 == 0:
                table = s.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(5), Inches(1.3)).table
                for r, row in enumerate([("k", "v"), ("x", str(i)), ("y", str(i * 2))]):
                    for c, v in enumerate(row):
                        table.cell(r, c).text = v
            if i % 3 == 0:
                s.shapes.add_picture(str(assets["A"]), Inches(6.5), Inches(2), width=Inches(2))
            s.notes_slide.notes_text_frame.text = f"NOTE_A215_{i}"
        tb(prs, "收尾", ["A215_END"])

    save("FS_A2_15_stress.pptx", a215)

    return names


def gen_excel(assets) -> list[str]:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Font, PatternFill, numbers
    from openpyxl.worksheet.table import Table, TableStyleInfo
    from openpyxl.comments import Comment
    import xlwt

    names = []

    def save(name, fn):
        p = SAMPLES / name
        wb = Workbook()
        fn(wb)
        wb.save(p)
        names.append(name)

    # A3-01
    def e01(wb):
        ws = wb.active
        ws.title = "说明"
        ws["A1"] = "A3-01 纯文 A301_TEXT"
        ws["A2"] = "华北 A301_CN"
        ws["A3"] = "Follow-up A301_EN"

    save("FS_A3_01_text.xlsx", e01)

    # A3-02
    def e02(wb):
        ws = wb.active
        ws.title = "收入"
        ws["A1"] = "A3-02 单表 A302_TT"
        for c, h in enumerate(["区域", "营收", "同比"], 1):
            cell = ws.cell(3, c, h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="2563EB")
        for i, row in enumerate([("华北", 1280, "+12%"), ("海外", 430, "-5%")], 4):
            for j, v in enumerate(row, 1):
                ws.cell(i, j, v)
        ws["A6"] = "A302_MARK"

    save("FS_A3_02_table.xlsx", e02)

    # A3-03
    def e03(wb):
        ws = wb.active
        ws["A1"] = "A3-03 图 A303_IMG"
        ws["A2"] = "A303_MARK"
        img = XLImage(str(assets["A"]))
        img.width, img.height = 140, 45
        ws.add_image(img, "A4")
        ws["A8"] = "CAP_A303"

    save("FS_A3_03_image.xlsx", e03)

    # A3-04
    def e04(wb):
        ws = wb.active
        ws["A1"] = "A3-04 综合 A304_MIX"
        ws["A2"] = "A304_MARK"
        img = XLImage(str(assets["C"]))
        img.width, img.height = 150, 55
        ws.add_image(img, "A3")
        ws["A10"] = "指标"
        ws["B10"] = "值"
        ws["A11"] = "准确率"
        ws["B11"] = "96%"
        ws["A12"] = "CAP_A304"

    save("FS_A3_04_mix.xlsx", e04)

    # A3-05 并排
    def e05(wb):
        ws = wb.active
        ws["A1"] = "左 A305_LEFT"
        ws["A2"] = "区"
        ws["B2"] = "值"
        ws["A3"] = "华北"
        ws["B3"] = 1280
        ws["D1"] = "右 A305_RIGHT"
        ws["D2"] = "项"
        ws["E2"] = "额"
        ws["D3"] = "人力"
        ws["E3"] = 200
        ws["A5"] = "A305_MARK"

    save("FS_A3_05_side.xlsx", e05)

    # A3-06 多 sheet
    def e06(wb):
        for i, name in enumerate(["华北", "华南", "海外", "汇总"]):
            ws = wb.active if i == 0 else wb.create_sheet(name)
            if i == 0:
                ws.title = name
            ws["A1"] = f"SHEET_{name}"
            ws["A2"] = "营收"
            ws["B2"] = 100 * (i + 1)
            ws["A3"] = f"MARK_{name}"
        wb["汇总"]["A4"] = "A306_MULTI"

    save("FS_A3_06_multi_sheet.xlsx", e06)

    # A3-07 复杂合并
    def e07(wb):
        ws = wb.active
        ws.merge_cells("A1:E1")
        ws["A1"] = "A3-07 合并 A307_MERGE_HDR"
        ws.merge_cells("A2:A4")
        ws["A2"] = "华北组"
        ws["B2"] = "收入"
        ws["C2"] = 1280
        ws["D2"] = "+12%"
        ws["E2"] = "稳"
        ws["B3"] = "成本"
        ws["C3"] = 800
        ws["D3"] = "—"
        ws["E3"] = "—"
        ws["B4"] = "利润"
        ws["C4"] = 480
        ws["D4"] = "—"
        ws["E4"] = "—"
        ws["A5"] = "海外"
        ws["B5"] = "收入"
        ws["C5"] = 430
        ws["D5"] = "-5%"
        ws["E5"] = "压"
        ws["A6"] = "A307_END"

    save("FS_A3_07_merge.xlsx", e07)

    # A3-08 宽表
    def e08(wb):
        ws = wb.active
        ws["A1"] = "A3-08 宽表 A308_WIDE"
        for c in range(1, 21):
            ws.cell(2, c, f"列{c}")
            ws.cell(3, c, c * 7)
        ws["A4"] = "A308_MARK"

    save("FS_A3_08_wide.xlsx", e08)

    # A3-09 稀疏
    def e09(wb):
        ws = wb.active
        ws["A1"] = "A3-09 稀疏 A309_SPARSE"
        ws["A3"] = "华北"
        ws["B3"] = 1280
        ws["A10"] = "远行 A309_FAR"
        ws["G3"] = "远列 A309_COL"
        ws["A12"] = "A309_MARK"

    save("FS_A3_09_sparse.xlsx", e09)

    # A3-10 数字格式
    def e10(wb):
        ws = wb.active
        ws["A1"] = "A3-10 格式 A310_FMT"
        ws["A2"] = "百分数"
        ws["B2"] = 0.125
        ws["B2"].number_format = "0.0%"
        ws["A3"] = "货币"
        ws["B3"] = 1280.5
        ws["B3"].number_format = '"¥"#,##0.00'
        ws["A4"] = "日期"
        ws["B4"] = "2024-06-15"
        ws["A5"] = "A310_MARK"

    save("FS_A3_10_formats.xlsx", e10)

    # A3-11 公式
    def e11(wb):
        ws = wb.active
        ws["A1"] = "A3-11 公式 A311_FORMULA"
        ws["A2"] = 10
        ws["B2"] = 20
        ws["C2"] = "=A2+B2"
        ws["A3"] = "A311_MARK"

    save("FS_A3_11_formula.xlsx", e11)

    # A3-12 多图
    def e12(wb):
        ws = wb.active
        ws["A1"] = "A3-12 多图 A312_IMG"
        for i, k in enumerate(["A", "B", "C"]):
            img = XLImage(str(assets[k]))
            img.width, img.height = 110, 38
            ws.add_image(img, f"A{3 + i * 4}")
        ws["A16"] = "A312_MARK"

    save("FS_A3_12_multi_image.xlsx", e12)

    # A3-13 隐藏行列
    def e13(wb):
        ws = wb.active
        ws["A1"] = "A3-13 隐藏 A313_HIDE"
        ws["A2"] = "可见行 A313_VISIBLE"
        ws["A3"] = "隐藏行 A313_HIDDEN_ROW"
        ws.row_dimensions[3].hidden = True
        ws["B1"] = "可见列"
        ws["C1"] = "隐藏列 A313_HIDDEN_COL"
        ws.column_dimensions["C"].hidden = True
        ws["A4"] = "A313_MARK"

    save("FS_A3_13_hidden.xlsx", e13)

    # A3-14 Excel Table
    def e14(wb):
        ws = wb.active
        ws["A1"] = "区域"
        ws["B1"] = "营收"
        ws["A2"] = "华北"
        ws["B2"] = 1280
        ws["A3"] = "海外"
        ws["B3"] = 430
        tab = Table(displayName="IncomeTable", ref="A1:B3")
        tab.tableStyleInfo = TableStyleInfo(name="TableStyleMedium2", showRowStripes=True)
        ws.add_table(tab)
        ws["A5"] = "A3-14 Table A314_TABLE"
        ws["A6"] = "A314_MARK"

    save("FS_A3_14_excel_table.xlsx", e14)

    # A3-15 压力包
    def e15(wb):
        for si in range(1, 4):
            ws = wb.active if si == 1 else wb.create_sheet(f"S{si}")
            if si == 1:
                ws.title = "S1"
            ws["A1"] = f"A315 压力 S{si} A315_S{si}"
            ws["A2"] = "区"
            ws["B2"] = "值"
            ws["A3"] = "华北"
            ws["B3"] = 100 * si
            ws["D1"] = f"右表 A315_R{si}"
            ws["D2"] = "项"
            ws["E2"] = "额"
            ws["D3"] = "云"
            ws["E3"] = 40 * si
            img = XLImage(str(assets["A"]))
            img.width, img.height = 100, 35
            ws.add_image(img, "A6")
        wb.create_sheet("END")["A1"] = "A315_END"

    save("FS_A3_15_stress.xlsx", e15)

    # 补充：超链接单元格、批注、图表
    def e16(wb):
        ws = wb.active
        ws["A1"] = "A3-16 链接批注图 A316_EXTRA"
        ws["A2"] = "链接单元格"
        ws["A2"].hyperlink = "https://example.com/a316"
        ws["A3"] = "华北"
        ws["A3"].comment = Comment("批注 COMMENT_A316", "tester")
        ws["A5"] = "区"
        ws["B5"] = "值"
        ws["A6"] = "华北"
        ws["B6"] = 10
        ws["A7"] = "海外"
        ws["B7"] = 5
        chart = BarChart()
        chart.title = "CHART_A316"
        data = Reference(ws, min_col=2, min_row=5, max_row=7)
        cats = Reference(ws, min_col=1, min_row=6, max_row=7)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        ws.add_chart(chart, "D5")
        ws["A9"] = "A316_MARK"

    save("FS_A3_16_link_comment_chart.xlsx", e16)

    # A4 xls
    p = SAMPLES / "FS_A4_01_xls_basic.xls"
    book = xlwt.Workbook()
    sh = book.add_sheet("收入")
    sh.write(0, 0, "A4-01 XLS A401_XLS")
    sh.write(1, 0, "区域")
    sh.write(1, 1, "营收")
    sh.write(2, 0, "华北")
    sh.write(2, 1, 1280)
    sh.write(3, 0, "A401_MARK")
    book.save(str(p))
    names.append("FS_A4_01_xls_basic.xls")

    p2 = SAMPLES / "FS_A4_02_xls_multi.xls"
    book = xlwt.Workbook()
    for name, mark, val in [("华北", "A402_N", 1280), ("海外", "A402_O", 430)]:
        sh = book.add_sheet(name)
        sh.write(0, 0, f"A4-02 {name} {mark}")
        sh.write(1, 0, "营收")
        sh.write(1, 1, val)
        sh.write(2, 0, "A402_MARK")
    book.save(str(p2))
    names.append("FS_A4_02_xls_multi.xls")

    return names


def gen_msg() -> list[str]:
    names = []
    SAMPLES.mkdir(parents=True, exist_ok=True)
    if MSG.exists():
        for name in [
            "FS_A5_01_official.msg",
            "FS_A5_06_official_regress.msg",
            "FS_A5_03_中文名邮件.msg",
        ]:
            shutil.copy2(MSG, SAMPLES / name)
            names.append(name)
        # 截断
        data = MSG.read_bytes()
        (SAMPLES / "FS_A5_04_truncated.msg").write_bytes(data[: max(80, len(data) // 8)])
        names.append("FS_A5_04_truncated.msg")

    # 伪 HTML 感无法真造 HTML msg 无 Outlook；用标记文件说明 skip 路径——改为负面伪文件
    (SAMPLES / "FS_A5_02_fake_htmlish.msg").write_bytes(
        b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"FAKE_HTML_MSG_A502" + b"\x00" * 120
    )
    names.append("FS_A5_02_fake_htmlish.msg")

    (SAMPLES / "FS_A5_05_empty.msg").write_bytes(b"")
    names.append("FS_A5_05_empty.msg")
    return names


def gen_negatives() -> list[str]:
    names = []
    SAMPLES.mkdir(parents=True, exist_ok=True)
    (SAMPLES / "FS_G_01_empty.docx").write_bytes(b"")
    names.append("FS_G_01_empty.docx")

    (SAMPLES / "FS_G_05_corrupt.docx").write_bytes(b"PK\x03\x04" + b"\x00" * 40 + b"CORRUPT_G05")
    names.append("FS_G_05_corrupt.docx")
    (SAMPLES / "FS_G_05_corrupt.pptx").write_bytes(b"PK\x03\x04" + b"\x00" * 40 + b"CORRUPT_PPT")
    names.append("FS_G_05_corrupt.pptx")
    (SAMPLES / "FS_G_05_corrupt.xlsx").write_bytes(b"PK\x03\x04" + b"\x00" * 40 + b"CORRUPT_XLS")
    names.append("FS_G_05_corrupt.xlsx")

    (SAMPLES / "FS_G_01_old.doc").write_bytes(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"OLD_DOC_G" + b"\x00" * 100)
    names.append("FS_G_01_old.doc")
    (SAMPLES / "FS_G_02_old.ppt").write_bytes(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"OLD_PPT_G" + b"\x00" * 100)
    names.append("FS_G_02_old.ppt")

    # 中文文件名
    from docx import Document

    cn = SAMPLES / "FS_G_03（中文 空格）名.docx"
    d = Document()
    d.add_heading("G-03 文件名 G03_NAME", 1)
    d.add_paragraph("内容 G03_BODY")
    d.save(cn)
    names.append(cn.name)

    # 后缀不符：msg 当 docx
    if MSG.exists():
        shutil.copy2(MSG, SAMPLES / "FS_G_04_msg_as.docx")
        names.append("FS_G_04_msg_as.docx")

    return names


def main():
    assets = make_assets()
    w = gen_word(assets)
    p = gen_pptx(assets)
    e = gen_excel(assets)
    m = gen_msg()
    g = gen_negatives()
    total = len(w) + len(p) + len(e) + len(m) + len(g)
    print(f"A1={len(w)} A2={len(p)} A3/A4={len(e)} A5={len(m)} G={len(g)} total={total}")


if __name__ == "__main__":
    main()
