"""Office 加深加压：Word / PPT / Excel / Outlook 穷举样例生成。"""

from __future__ import annotations

import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parent
SAMPLES = ROOT / "samples"
ASSETS = ROOT / "assets"
PREV_MSG = Path(__file__).resolve().parents[1] / "full_matrix" / "phase_retry" / "test_outlook_msg.msg"


def _font(size: int = 20):
    for name in (r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\arial.ttf"):
        if Path(name).exists():
            try:
                return ImageFont.truetype(name, size)
            except OSError:
                pass
    return ImageFont.load_default()


def make_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    out = {}
    for key, color, size in [
        ("N", (37, 99, 235), (200, 70)),
        ("S", (5, 150, 105), (200, 70)),
        ("C", (220, 38, 38), (240, 100)),
        ("L", (80, 80, 80), (100, 50)),
        ("X", (124, 58, 237), (180, 70)),
    ]:
        p = ASSETS / f"img_{key}.png"
        im = Image.new("RGB", size, color)
        dr = ImageDraw.Draw(im)
        dr.text((10, size[1] // 2 - 10), f"IMG_{key}_MARK", fill=(255, 255, 255), font=_font(18))
        im.save(p)
        out[key] = p
    return out


# ---------- Word ----------
def gen_word(assets: dict[str, Path]) -> list[Path]:
    from docx import Document
    from docx.shared import Inches, RGBColor
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn

    paths = []
    SAMPLES.mkdir(parents=True, exist_ok=True)

    def save(name, fn):
        p = SAMPLES / name
        doc = Document()
        fn(doc)
        doc.save(p)
        paths.append(p)

    def add_table(doc, rows):
        t = doc.add_table(rows=len(rows), cols=len(rows[0]))
        t.style = "Table Grid"
        for i, row in enumerate(rows):
            for j, v in enumerate(row):
                t.rows[i].cells[j].text = str(v)
        return t

    save("WD01_text_only.docx", lambda d: (
        d.add_heading("WD01 纯文字 WD_TEXT", 1),
        d.add_heading("概述", 2),
        d.add_paragraph("华北增长12%。海外-5%。WD_BODY_CN"),
        d.add_paragraph("Action item follow-up. WD_BODY_EN"),
        d.add_paragraph("要点A", style="List Bullet"),
        d.add_paragraph("要点B", style="List Bullet"),
        d.add_paragraph("步骤1", style="List Number"),
    ) and None)

    def wd02(d):
        d.add_heading("WD02 文+表 WD_TT", 1)
        d.add_paragraph("见表 WD_TT_MARK")
        add_table(d, [("区域", "营收", "同比"), ("华北", "1280", "+12%"), ("海外", "430", "-5%")])

    save("WD02_text_table.docx", wd02)

    def wd03(d):
        d.add_heading("WD03 文+图 WD_TI", 1)
        d.add_paragraph("见图 WD_TI_MARK")
        d.add_picture(str(assets["N"]), width=Inches(2))
        d.add_paragraph("图注 CAP_WD_N")

    save("WD03_text_image.docx", wd03)

    def wd04(d):
        d.add_heading("WD04 文图表 WD_MIX", 1)
        d.add_paragraph("综合 WD_MIX_MARK")
        d.add_picture(str(assets["C"]), width=Inches(2.2))
        d.add_paragraph("图注 CAP_WD_C")
        add_table(d, [("指标", "值"), ("准确率", "96%"), ("召回率", "91%")])
        d.add_paragraph("结尾 WD_MIX_END")

    save("WD04_text_image_table.docx", wd04)

    def wd05(d):
        d.add_heading("WD05 多表 WD_MT", 1)
        d.add_paragraph("表1 WD_MT1")
        add_table(d, [("区", "值"), ("华北", "1280"), ("华南", "960")])
        d.add_paragraph("表2 WD_MT2")
        add_table(d, [("项", "额"), ("人力", "200"), ("云", "80")])
        d.add_paragraph("表3 WD_MT3")
        add_table(d, [("风险", "等级"), ("海外", "高"), ("合规", "中")])

    save("WD05_multi_table.docx", wd05)

    def wd06(d):
        d.add_heading("WD06 多图 WD_MI", 1)
        d.add_paragraph("多图 WD_MI_MARK")
        for k, cap in [("N", "CAP_WD_N1"), ("S", "CAP_WD_S1"), ("L", "CAP_WD_L1"), ("X", "CAP_WD_X1")]:
            d.add_picture(str(assets[k]), width=Inches(1.6))
            d.add_paragraph(f"图注 {cap}")

    save("WD06_multi_image.docx", wd06)

    def wd07(d):
        d.add_heading("WD07 合并表 WD_MERGE", 1)
        t = add_table(d, [["合计区 MERGE_HDR", "", ""], ["华北", "1280", "+12%"], ["海外", "430", "-5%"]])
        t.cell(0, 0).merge(t.cell(0, 2))

    save("WD07_merged_table.docx", wd07)

    def wd08(d):
        d.add_heading("WD08 多节长文 WD_LONG", 1)
        for i in range(1, 6):
            d.add_heading(f"第{i}章 CH{i}", 2)
            d.add_heading(f"小节 {i}.1 SEC{i}", 3)
            d.add_paragraph(f"段落 {i} WD_LONG_P{i} mixed 中英")

    save("WD08_multi_section.docx", wd08)

    def wd09(d):
        d.add_heading("WD09 中英混排 WD_MIXLANG", 1)
        d.add_paragraph("Q2 Revenue 季度营收：North China（华北）+12%。WD_LANG_MARK")
        r = d.add_paragraph().add_run("Red alert WD_RED")
        r.bold = True
        r.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)

    save("WD09_mixed_lang.docx", wd09)

    def wd10(d):
        d.add_heading("WD10 压力包 WD_STRESS", 1)
        d.add_paragraph("压力开篇 WD_STRESS_MARK")
        for i in range(1, 4):
            d.add_heading(f"块{i}", 2)
            d.add_picture(str(assets["N" if i % 2 else "S"]), width=Inches(1.5))
            d.add_paragraph(f"图注 CAP_ST_{i}")
            add_table(d, [("k", "v"), (f"row{i}a", str(100 * i)), (f"row{i}b", str(200 * i))])
        d.add_paragraph("压力结尾 WD_STRESS_END")

    save("WD10_stress_pack.docx", wd10)

    def wd11(d):
        sec = d.sections[0]
        sec.header.paragraphs[0].text = "HEADER_WD"
        sec.footer.paragraphs[0].text = "FOOTER_WD"
        d.add_heading("WD11 页眉页脚 WD_HF", 1)
        d.add_paragraph("正文 BODY_WD_HF")

    save("WD11_header_footer.docx", wd11)

    def wd12(d):
        d.add_heading("WD12 超链接 WD_LINK", 1)
        p = d.add_paragraph()
        rid = p.part.relate_to(
            "https://example.com/wd",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
        hl = OxmlElement("w:hyperlink")
        hl.set(qn("r:id"), rid)
        nr = OxmlElement("w:r")
        te = OxmlElement("w:t")
        te.text = "链接文字 LINK_WD"
        nr.append(te)
        hl.append(nr)
        p._p.append(hl)
        d.add_paragraph("链接后 WD_LINK_BODY")

    save("WD12_hyperlink.docx", wd12)

    # 负面：空几乎
    def wd13(d):
        d.add_paragraph("")

    save("WD13_almost_empty.docx", wd13)

    return paths


# ---------- PPT ----------
def gen_pptx(assets: dict[str, Path]) -> list[Path]:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    paths = []

    def save(name, fn):
        p = SAMPLES / name
        prs = Presentation()
        fn(prs)
        prs.save(p)
        paths.append(p)

    def add_title_body(prs, title, body, notes=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.text = body[0] if body else ""
        for line in body[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return slide

    def ppt01(prs):
        for i in range(1, 6):
            add_title_body(prs, f"PPT纯文页{i} PPT_T{i}", [f"正文 PPT_TEXT_{i}", "中英 mixed content"])

    save("PP01_text_only.pptx", ppt01)

    def ppt02(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "PPT文+表 PPT_TT"
        rows, cols = 4, 3
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(2)).table
        data = [("区域", "营收", "同比"), ("华北", "1280", "+12%"), ("华南", "960", "0%"), ("海外", "430", "-5%")]
        for r, row in enumerate(data):
            for c, v in enumerate(row):
                table.cell(r, c).text = v
        # 说明页
        add_title_body(prs, "说明", ["见表 PPT_TT_MARK"])

    save("PP02_text_table.pptx", ppt02)

    def ppt03(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "PPT文+图 PPT_TI"
        slide.shapes.add_picture(str(assets["N"]), Inches(1), Inches(2), width=Inches(3))
        add_title_body(prs, "图注页", ["CAPTION_PPT_N", "PPT_TI_MARK"])

    save("PP03_text_image.pptx", ppt03)

    def ppt04(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "PPT文图表 PPT_MIX"
        slide.shapes.add_picture(str(assets["C"]), Inches(0.5), Inches(1.8), width=Inches(2.5))
        table = slide.shapes.add_table(3, 2, Inches(4), Inches(2), Inches(5), Inches(1.5)).table
        for r, row in enumerate([("指标", "值"), ("准确率", "96%"), ("召回率", "91%")]):
            for c, v in enumerate(row):
                table.cell(r, c).text = v
        add_title_body(prs, "综合说明", ["PPT_MIX_MARK", "CAP_PPT_C"])

    save("PP04_text_image_table.pptx", ppt04)

    def ppt05(prs):
        for i, title in enumerate(["表A PPT_MT1", "表B PPT_MT2", "表C PPT_MT3"], 1):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            table = slide.shapes.add_table(3, 2, Inches(1.5), Inches(2), Inches(6), Inches(1.5)).table
            for r, row in enumerate([("k", "v"), (f"a{i}", str(i * 10)), (f"b{i}", str(i * 20))]):
                for c, v in enumerate(row):
                    table.cell(r, c).text = v

    save("PP05_multi_table.pptx", ppt05)

    def ppt06(prs):
        for k, cap in [("N", "CAP_PPT_N1"), ("S", "CAP_PPT_S1"), ("X", "CAP_PPT_X1")]:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"多图 {cap}"
            slide.shapes.add_picture(str(assets[k]), Inches(1), Inches(2), width=Inches(3))
            add_title_body(prs, "说明", [cap, "PPT_MI_MARK"])

    save("PP06_multi_image.pptx", ppt06)

    def ppt07(prs):
        for i in range(1, 6):
            add_title_body(
                prs,
                f"多节页{i} PPT_SEC{i}",
                [f"内容 PPT_SECTION_{i}", "中英 Section content"],
                notes=f"备注 NOTE_PPT_{i}",
            )

    save("PP07_multi_section_notes.pptx", ppt07)

    def ppt08(prs):
        slide = add_title_body(prs, "颜色样式 PPT_STYLE", ["红标题样式", "PPT_STYLE_MARK"])
        p = slide.shapes.title.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        p.font.size = Pt(36)
        body = slide.placeholders[1].text_frame
        p2 = body.add_paragraph()
        p2.text = "二级要点"
        p2.level = 1
        p2.font.color.rgb = RGBColor(0x05, 0x96, 0x69)

    save("PP08_styles_list.pptx", ppt08)

    def ppt09(prs):
        # 正常页 + 隐藏页
        add_title_body(prs, "可见页 PPT_VISIBLE", ["VISIBLE_MARK"])
        slide = add_title_body(prs, "隐藏页 PPT_HIDDEN", ["HIDDEN_MARK"])
        slide._element.set("show", "0")  # hide

    save("PP09_hidden_slide.pptx", ppt09)

    def ppt10(prs):
        add_title_body(prs, "中英混排 PPT_LANG", ["Q2 Revenue 华北 +12%", "Overseas 海外 -5% PPT_LANG_MARK"])

    save("PP10_mixed_lang.pptx", ppt10)

    def ppt11(prs):
        for i in range(1, 11):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"压力{i} PPT_STRESS_{i}"
            if i % 3 == 0:
                slide.shapes.add_picture(str(assets["N"]), Inches(0.8), Inches(1.8), width=Inches(2))
            if i % 2 == 0:
                table = slide.shapes.add_table(3, 2, Inches(4), Inches(2), Inches(5), Inches(1.4)).table
                for r, row in enumerate([("k", "v"), ("x", str(i)), ("y", str(i * 2))]):
                    for c, v in enumerate(row):
                        table.cell(r, c).text = v
            slide.notes_slide.notes_text_frame.text = f"NOTE_STRESS_{i}"
        add_title_body(prs, "压力收尾", ["PPT_STRESS_END"])

    save("PP11_stress_pack.pptx", ppt11)

    def ppt12(prs):
        # 几乎空
        prs.slides.add_slide(prs.slide_layouts[6])

    save("PP12_blankish.pptx", ppt12)

    return paths


# ---------- Excel ----------
def gen_excel(assets: dict[str, Path]) -> list[Path]:
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Font, PatternFill
    import xlwt

    paths = []

    def save_xlsx(name, fn):
        p = SAMPLES / name
        wb = Workbook()
        fn(wb)
        wb.save(p)
        paths.append(p)

    save_xlsx("EX01_text_only.xlsx", lambda wb: (
        setattr(wb.active, "title", "说明"),
        wb.active.__setitem__("A1", "EX01 纯文字 EX_TEXT"),
        wb.active.__setitem__("A2", "华北增长 EX_BODY_CN"),
        wb.active.__setitem__("A3", "Follow up EX_BODY_EN"),
    ) and None)

    def ex02(wb):
        ws = wb.active
        ws.title = "收入"
        ws["A1"] = "EX02 文+表 EX_TT"
        ws["A2"] = "EX_TT_MARK"
        for c, h in enumerate(["区域", "营收", "同比"], 1):
            cell = ws.cell(4, c, h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="2563EB")
        for i, row in enumerate([("华北", 1280, "+12%"), ("海外", 430, "-5%")], 5):
            for j, v in enumerate(row, 1):
                ws.cell(i, j, v)

    save_xlsx("EX02_text_table.xlsx", ex02)

    def ex03(wb):
        ws = wb.active
        ws.title = "图"
        ws["A1"] = "EX03 文+图 EX_TI"
        ws["A2"] = "EX_TI_MARK"
        img = XLImage(str(assets["N"]))
        img.width, img.height = 160, 55
        ws.add_image(img, "A4")
        ws["A8"] = "CAP_EX_N"

    save_xlsx("EX03_text_image.xlsx", ex03)

    def ex04(wb):
        ws = wb.active
        ws.title = "综合"
        ws["A1"] = "EX04 文图表 EX_MIX"
        ws["A2"] = "EX_MIX_MARK"
        img = XLImage(str(assets["C"]))
        img.width, img.height = 180, 70
        ws.add_image(img, "A3")
        ws["A10"] = "指标"
        ws["B10"] = "值"
        ws["A11"] = "准确率"
        ws["B11"] = "96%"
        ws["A12"] = "CAP_EX_C"

    save_xlsx("EX04_text_image_table.xlsx", ex04)

    def ex05(wb):
        ws = wb.active
        ws.title = "并排"
        ws["A1"] = "左表 EX_MT_LEFT"
        ws["A2"] = "区"
        ws["B2"] = "值"
        ws["A3"] = "华北"
        ws["B3"] = 1280
        ws["D1"] = "右表 EX_MT_RIGHT"
        ws["D2"] = "项"
        ws["E2"] = "额"
        ws["D3"] = "人力"
        ws["E3"] = 200
        ws["G1"] = "第三表 EX_MT_THIRD"
        ws["G2"] = "风险"
        ws["H2"] = "级"
        ws["G3"] = "海外"
        ws["H3"] = "高"

    save_xlsx("EX05_multi_table_side.xlsx", ex05)

    def ex06(wb):
        ws = wb.active
        ws.title = "华北"
        ws["A1"] = "SHEET_N"
        ws["A2"] = "营收"
        ws["B2"] = 1280
        wb.create_sheet("华南")["A1"] = "SHEET_S"
        wb["华南"]["A2"] = "营收"
        wb["华南"]["B2"] = 960
        wb.create_sheet("海外")["A1"] = "SHEET_O"
        wb["海外"]["A2"] = "营收"
        wb["海外"]["B2"] = 430
        wb.create_sheet("汇总")["A1"] = "SHEET_SUM"
        wb["汇总"]["A2"] = "合计"
        wb["汇总"]["B2"] = 2670

    save_xlsx("EX06_multi_sheet.xlsx", ex06)

    def ex07(wb):
        ws = wb.active
        ws.title = "合并"
        ws.merge_cells("A1:D1")
        ws["A1"] = "EX07 合并 EX_MERGE_HDR"
        ws.merge_cells("A2:A4")
        ws["A2"] = "华北组"
        ws["B2"] = "收入"
        ws["C2"] = 1280
        ws["D2"] = "+12%"
        ws["B3"] = "成本"
        ws["C3"] = 800
        ws["D3"] = "—"
        ws["B4"] = "利润"
        ws["C4"] = 480
        ws["D4"] = "—"
        ws["A5"] = "海外"
        ws["B5"] = "收入"
        ws["C5"] = 430
        ws["D5"] = "-5%"

    save_xlsx("EX07_merged_nested.xlsx", ex07)

    def ex08(wb):
        ws = wb.active
        ws.title = "混排"
        ws["A1"] = "EX08 Q2 Revenue 季度 EX_LANG"
        ws["A2"] = "North China（华北）1280"
        ws["A3"] = "Overseas（海外）430 EX_LANG_MARK"

    save_xlsx("EX08_mixed_lang.xlsx", ex08)

    def ex09(wb):
        ws = wb.active
        ws.title = "宽表"
        ws["A1"] = "EX09 宽表 EX_WIDE"
        for i in range(1, 21):
            ws.cell(2, i, f"列{i}")
            ws.cell(3, i, i * 11)
        ws["A4"] = "EX_WIDE_MARK"

    save_xlsx("EX09_wide_table.xlsx", ex09)

    def ex10(wb):
        for si in range(1, 4):
            ws = wb.active if si == 1 else wb.create_sheet(f"S{si}")
            if si == 1:
                ws.title = "S1"
            ws["A1"] = f"EX10 压力 Sheet{si} EX_STRESS_{si}"
            ws["A2"] = "区"
            ws["B2"] = "值"
            ws["A3"] = "华北"
            ws["B3"] = 100 * si
            ws["D1"] = f"右表 R{si}"
            ws["D2"] = "项"
            ws["E2"] = "额"
            ws["D3"] = "云"
            ws["E3"] = 50 * si
            try:
                img = XLImage(str(assets["N"]))
                img.width, img.height = 120, 40
                ws.add_image(img, "A6")
            except Exception:
                pass
        wb.create_sheet("END")["A1"] = "EX_STRESS_END"

    save_xlsx("EX10_stress_pack.xlsx", ex10)

    def ex11(wb):
        ws = wb.active
        ws.title = "公式格式"
        ws["A1"] = "EX11 EX_FMT"
        ws["A2"] = 0.12
        ws["A2"].number_format = "0%"
        ws["B2"] = 1280
        ws["B2"].number_format = '"¥"#,##0'
        ws["C2"] = "=B2*A2"
        ws["A3"] = "EX_FMT_MARK"

    save_xlsx("EX11_formula_format.xlsx", ex11)

    # 负面：几乎空
    def ex12(wb):
        wb.active["A1"] = ""

    save_xlsx("EX12_emptyish.xlsx", ex12)

    # xls
    p = SAMPLES / "EX13_xls_basic.xls"
    book = xlwt.Workbook()
    sh = book.add_sheet("收入")
    sh.write(0, 0, "EX13 XLS EX_XLS")
    sh.write(1, 0, "区域")
    sh.write(1, 1, "营收")
    sh.write(2, 0, "华北")
    sh.write(2, 1, 1280)
    sh.write(3, 0, "EX_XLS_MARK")
    book.save(str(p))
    paths.append(p)

    return paths


# ---------- Outlook MSG ----------
def gen_msg() -> list[Path]:
    paths = []
    SAMPLES.mkdir(parents=True, exist_ok=True)
    # 官方真实样例
    if PREV_MSG.exists():
        dest = SAMPLES / "MSG01_official.msg"
        shutil.copy2(PREV_MSG, dest)
        paths.append(dest)

    # 再下一份公开样例（若网络可用）
    try:
        import urllib.request

        url = "https://raw.githubusercontent.com/microsoft/markitdown/main/packages/markitdown/tests/test_files/test_outlook_msg.msg"
        dest2 = SAMPLES / "MSG02_official_dl.msg"
        urllib.request.urlretrieve(url, dest2)
        paths.append(dest2)
    except Exception:
        pass

    # 伪 msg 负面
    bad = SAMPLES / "MSG03_fake_ole.msg"
    bad.write_bytes(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"\x00" * 128)
    paths.append(bad)

    return paths


def main():
    assets = make_assets()
    w = gen_word(assets)
    p = gen_pptx(assets)
    e = gen_excel(assets)
    m = gen_msg()
    print(f"word={len(w)} ppt={len(p)} excel={len(e)} msg={len(m)} total={len(w)+len(p)+len(e)+len(m)}")


if __name__ == "__main__":
    main()
