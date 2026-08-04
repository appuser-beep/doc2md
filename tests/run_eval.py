"""生成带字体颜色/字号的 Office 样例，并批量转换评估 LLM 友好度。"""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT))

SAMPLES = Path(__file__).resolve().parent / "samples"
OUTPUT = Path(__file__).resolve().parent / "output"
REPORT = Path(__file__).resolve().parent / "EVAL_REPORT.md"


def make_docx() -> Path:
    from docx import Document
    from docx.shared import Pt, RGBColor

    path = SAMPLES / "06_styled_report.docx"
    doc = Document()
    h = doc.add_heading("季度经营报告（Word）", level=1)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        run.font.size = Pt(28)

    h2 = doc.add_heading("收入概况", level=2)
    for run in h2.runs:
        run.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
        run.font.size = Pt(20)

    p = doc.add_paragraph()
    r = p.add_run("红色大号强调：华北增长 12%。")
    r.font.color.rgb = RGBColor(0xE1, 0x1D, 0x48)
    r.font.size = Pt(22)
    r.bold = True

    p2 = doc.add_paragraph()
    r2 = p2.add_run("蓝色小号说明：本段仅用于视觉排版测试。")
    r2.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)
    r2.font.size = Pt(9)

    doc.add_paragraph("普通正文：管理层要求复盘海外下滑原因。")
    doc.add_paragraph("重点客户跟进", style="List Bullet")
    doc.add_paragraph("渠道漏斗复盘", style="List Bullet")
    doc.add_paragraph("输出行动清单", style="List Number")

    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    headers = ("区域", "营收(万元)", "同比")
    rows = [
        ("华北", "1280", "+12%"),
        ("华南", "960", "0%"),
        ("海外", "430", "-5%"),
    ]
    for i, htxt in enumerate(headers):
        table.rows[0].cells[i].text = htxt
    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            table.rows[ri].cells[ci].text = val

    doc.save(path)
    return path


def make_xlsx() -> Path:
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    path = SAMPLES / "07_styled_sheet.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "收入"
    headers = ["区域", "营收(万元)", "同比", "备注"]
    data = [
        ["华北", 1280, "+12%", "重点增长区"],
        ["华南", 960, "0%", "持平"],
        ["海外", 430, "-5%", "需复盘"],
    ]
    for col, h in enumerate(headers, 1):
        cell = ws.cell(1, col, h)
        cell.font = Font(color="FFFFFF", bold=True, size=14, name="微软雅黑")
        cell.fill = PatternFill("solid", fgColor="2563EB")
        cell.alignment = Alignment(horizontal="center")
    for r, row in enumerate(data, 2):
        for c, val in enumerate(row, 1):
            cell = ws.cell(r, c, val)
            cell.font = Font(size=11, name="Arial")
            if c == 3 and isinstance(val, str) and val.startswith("+"):
                cell.font = Font(color="059669", size=11, bold=True)
            if c == 3 and isinstance(val, str) and val.startswith("-"):
                cell.font = Font(color="DC2626", size=11, bold=True)
    wb.save(path)
    return path


def make_pptx() -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    path = SAMPLES / "08_styled_deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "季度经营报告"
    title_tf = slide.shapes.title.text_frame.paragraphs[0]
    title_tf.font.size = Pt(36)
    title_tf.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)

    body = slide.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = "华北增长 12%"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
    p2 = body.add_paragraph()
    p2.text = "海外下滑 5%（需复盘）"
    p2.level = 0
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
    p3 = body.add_paragraph()
    p3.text = "行动：跟进重点客户 / 复盘漏斗"
    p3.level = 1
    p3.font.size = Pt(12)
    p3.font.color.rgb = RGBColor(0x25, 0x63, 0xEB)

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "数据表"
    rows, cols = 4, 3
    table = slide2.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(8), Inches(2)).table
    values = [
        ("区域", "营收", "同比"),
        ("华北", "1280", "+12%"),
        ("华南", "960", "0%"),
        ("海外", "430", "-5%"),
    ]
    for r, row in enumerate(values):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    prs.save(path)
    return path


def analyze_markdown(name: str, md: str) -> dict:
    """启发式评估：结构保留与 LLM 友好度（非视觉保真）。"""
    checks = {
        "has_heading": bool(re.search(r"(?m)^#{1,6}\s+\S", md)),
        "has_list": bool(re.search(r"(?m)^\s*([-*+]|\d+\.)\s+\S", md)),
        "has_table": "|" in md and re.search(r"(?m)^\|.+\|", md) is not None,
        "has_link": "](" in md or "<http" in md.lower(),
        "has_code": "```" in md or re.search(r"(?m)^ {4}\S", md) is not None,
        "has_key_semantics": any(
            k in md for k in ("华北", "海外", "1280", "12%", "复盘", "经营")
        ),
        # 负面：残留 CSS/内联样式/无意义色值堆砌
        # 检测真正的样式噪声（CSS 属性 / style=），不把正文里提到的字体名算进去
        "no_css_noise": not re.search(
            r"(font-size\s*:|color\s*:|style\s*=|rgb\(|rgba\()", md, re.I
        ),
        "no_html_heavy": md.count("<") < 8,
        "reasonable_length": 40 <= len(md.strip()) <= 20000,
    }

    # 针对不同输入的期望
    expect = {
        "01_styles_structure.html": [
            "has_heading",
            "has_list",
            "has_table",
            "has_key_semantics",
            "no_css_noise",
            "reasonable_length",
        ],
        "02_metrics.json": ["has_key_semantics", "no_css_noise", "reasonable_length"],
        "03_regions.csv": ["has_key_semantics", "reasonable_length"],
        "04_already_md.md": ["has_heading", "has_list", "has_key_semantics", "reasonable_length"],
        "05_plain.txt": ["has_key_semantics", "reasonable_length"],
        "06_styled_report.docx": [
            "has_heading",
            "has_list",
            "has_table",
            "has_key_semantics",
            "no_css_noise",
            "reasonable_length",
        ],
        "07_styled_sheet.xlsx": ["has_key_semantics", "reasonable_length"],
        "08_styled_deck.pptx": ["has_key_semantics", "no_css_noise", "reasonable_length"],
    }
    needed = expect.get(name, ["has_key_semantics", "reasonable_length", "no_css_noise"])
    passed = sum(1 for k in needed if checks.get(k))
    score = round(100 * passed / max(len(needed), 1))

    # LLM 友好分：结构信号 + 去装饰
    llm_signals = [
        checks["has_heading"] or checks["has_list"] or checks["has_table"],
        checks["has_key_semantics"],
        checks["no_css_noise"],
        checks["no_html_heavy"],
        checks["reasonable_length"],
    ]
    llm_score = round(100 * sum(1 for x in llm_signals if x) / len(llm_signals))

    return {
        "checks": checks,
        "needed": needed,
        "structure_score": score,
        "llm_score": llm_score,
        "chars": len(md),
        "lines": md.count("\n") + 1,
    }


def llm_readability_notes(md: str) -> list[str]:
    notes = []
    if re.search(r"(?m)^#{1,6}\s+", md):
        notes.append("标题层级清晰，利于模型分段理解。")
    else:
        notes.append("缺少 Markdown 标题，模型需自行推断章节。")
    if re.search(r"(?m)^\|.+\|", md):
        notes.append("表格以 Markdown 呈现，便于抽取数值。")
    if re.search(r"(?m)^\s*([-*+]|\d+\.)\s+", md):
        notes.append("列表结构保留，适合作为行动项/要点输入。")
    if re.search(r"font-size\s*:|color\s*:|#e11d48", md, re.I):
        notes.append("仍残留样式噪声，可能浪费 token、干扰语义。")
    else:
        notes.append("未保留字号/颜色等视觉样式（对 LLM 通常是优点）。")
    if len(md) > 8000:
        notes.append("篇幅较长，喂给模型时建议再切片。")
    return notes


def main() -> None:
    OUTPUT.mkdir(parents=True, exist_ok=True)
    SAMPLES.mkdir(parents=True, exist_ok=True)

    print("生成 Office 样例…")
    make_docx()
    make_xlsx()
    make_pptx()

    from converter import convert_path

    results = []
    files = sorted(SAMPLES.iterdir())
    for path in files:
        if not path.is_file():
            continue
        print(f"转换：{path.name}")
        try:
            md = convert_path(str(path))
            out = OUTPUT / f"{path.stem}.md"
            out.write_text(md, encoding="utf-8")
            analysis = analyze_markdown(path.name, md)
            analysis["notes"] = llm_readability_notes(md)
            analysis["ok"] = True
            analysis["error"] = None
            analysis["preview"] = md[:1200]
            results.append({"file": path.name, "out": out.name, **analysis})
            print(
                f"  OK  structure={analysis['structure_score']}  "
                f"llm={analysis['llm_score']}  chars={analysis['chars']}"
            )
        except Exception as exc:  # noqa: BLE001
            results.append(
                {
                    "file": path.name,
                    "ok": False,
                    "error": str(exc),
                    "structure_score": 0,
                    "llm_score": 0,
                    "notes": [],
                    "preview": "",
                    "checks": {},
                    "needed": [],
                    "chars": 0,
                    "lines": 0,
                    "out": "",
                }
            )
            print(f"  FAIL  {exc}")

    # 写报告
    avg_llm = (
        round(sum(r["llm_score"] for r in results if r["ok"]) / max(sum(1 for r in results if r["ok"]), 1))
    )
    avg_struct = (
        round(
            sum(r["structure_score"] for r in results if r["ok"])
            / max(sum(1 for r in results if r["ok"]), 1)
        )
    )
    ok_n = sum(1 for r in results if r["ok"])

    lines = [
        "# 转换质量与 LLM 友好度测试报告",
        "",
        f"- 样例数：{len(results)}",
        f"- 成功转换：{ok_n}/{len(results)}",
        f"- 结构保留均分：{avg_struct}/100",
        f"- LLM 友好均分：{avg_llm}/100",
        "",
        "## 结论摘要",
        "",
        "1. **字体颜色 / 字号**：MarkItDown（及 Markdown 本身）通常**不会保留**视觉样式；这反而更适合大模型——减少噪声 token，聚焦语义。",
        "2. **便于 LLM 理解的部分**：标题、列表、表格、链接、关键数值与行动项大多能保留。",
        "3. **需要注意**：纯装饰信息（红字强调、黄底高亮）会丢失；若业务依赖“颜色=风险等级”，应在源文档用文字写明（如“风险：红色/下滑”），不要只靠颜色。",
        "",
        "## 逐项结果",
        "",
    ]

    for r in results:
        lines.append(f"### {r['file']}")
        if not r["ok"]:
            lines.append(f"- 状态：失败 — `{r['error']}`")
            lines.append("")
            continue
        lines.append(f"- 输出：`tests/output/{r['out']}`")
        lines.append(f"- 结构分：{r['structure_score']}/100 · LLM 友好分：{r['llm_score']}/100")
        lines.append(f"- 规模：{r['chars']} 字符 / {r['lines']} 行")
        lines.append("- 检查项：")
        for k in r["needed"]:
            mark = "✓" if r["checks"].get(k) else "✗"
            lines.append(f"  - {mark} `{k}`")
        lines.append("- 可读性备注：")
        for n in r["notes"]:
            lines.append(f"  - {n}")
        lines.append("")
        lines.append("<details><summary>预览（前 1200 字符）</summary>")
        lines.append("")
        lines.append("```markdown")
        lines.append(r["preview"])
        lines.append("```")
        lines.append("")
        lines.append("</details>")
        lines.append("")

    lines.extend(
        [
            "## 对工具使用的建议",
            "",
            "| 场景 | 建议 |",
            "|------|------|",
            "| 给大模型做摘要/问答 | 直接用本工具输出的 Markdown，效果通常足够好 |",
            "| 依赖颜色表达风险 | 在原文补充文字标签，勿只靠红/绿字 |",
            "| 复杂扫描件 PDF | 可能缺版式；可后续接 OCR/云服务 |",
            "| Token 控制 | 先转 Markdown 再切片，比直接塞 PDF 二进制更高效 |",
            "",
            "## 机器可读结果",
            "",
            "```json",
            json.dumps(
                [
                    {
                        "file": r["file"],
                        "ok": r["ok"],
                        "structure_score": r["structure_score"],
                        "llm_score": r["llm_score"],
                        "chars": r["chars"],
                        "error": r.get("error"),
                    }
                    for r in results
                ],
                ensure_ascii=False,
                indent=2,
            ),
            "```",
            "",
        ]
    )

    REPORT.write_text("\n".join(lines), encoding="utf-8")
    summary = {
        "ok": ok_n,
        "total": len(results),
        "avg_structure": avg_struct,
        "avg_llm": avg_llm,
        "report": str(REPORT),
    }
    print(json.dumps(summary, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
