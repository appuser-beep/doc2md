"""Office 穷举 Wave2（基础加强）+ Wave3（深入边界）。

覆盖：docx / pptx / xlsx / xls / msg
维度：样式列表 / 嵌套合并 / 中英特殊字符 / 长文多页 / 隐藏表 / 公式错误值 /
      文件名边界 / 损坏伪文件 / 老格式负面 / 压力体积
"""

from __future__ import annotations

import shutil
import zipfile
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parent
SAMPLES = ROOT / "samples"
ASSETS = ROOT / "assets"
PREV_MSG = Path(__file__).resolve().parents[1] / "full_matrix" / "phase_retry" / "test_outlook_msg.msg"
DEEP_MSG = Path(__file__).resolve().parents[1] / "office_deep" / "samples" / "MSG01_official.msg"


def _font(size: int = 18):
    for name in (r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\arial.ttf"):
        if Path(name).exists():
            try:
                return ImageFont.truetype(name, size)
            except OSError:
                pass
    return ImageFont.load_default()


def make_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    out = {}
    for key, color, size in [
        ("A", (37, 99, 235), (180, 60)),
        ("B", (5, 150, 105), (180, 60)),
        ("C", (220, 38, 38), (220, 90)),
        ("D", (124, 58, 237), (160, 60)),
        ("E", (30, 30, 30), (90, 45)),
    ]:
        p = ASSETS / f"ex_{key}.png"
        im = Image.new("RGB", size, color)
        dr = ImageDraw.Draw(im)
        dr.text((8, size[1] // 2 - 8), f"IMG_{key}", fill=(255, 255, 255), font=_font(16))
        im.save(p)
        out[key] = p
    return out


# ===================== Word =====================
def gen_word(assets: dict[str, Path]) -> list[str]:
    from docx import Document
    from docx.enum.text import WD_COLOR_INDEX
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor

    names: list[str] = []
    SAMPLES.mkdir(parents=True, exist_ok=True)

    def save(name: str, fn) -> None:
        p = SAMPLES / name
        doc = Document()
        fn(doc)
        doc.save(p)
        names.append(name)

    def table(doc, rows):
        t = doc.add_table(rows=len(rows), cols=len(rows[0]))
        t.style = "Table Grid"
        for i, row in enumerate(rows):
            for j, v in enumerate(row):
                t.rows[i].cells[j].text = str(v)
        return t

    # ---- Wave2 基础加强 ----
    def w20(d):
        d.add_heading("W20 样式基础 W20_STYLES", 1)
        p = d.add_paragraph()
        r = p.add_run("粗体 W20_BOLD ")
        r.bold = True
        r = p.add_run("斜体 W20_ITALIC ")
        r.italic = True
        r = p.add_run("下划线 W20_UL ")
        r.underline = True
        r = p.add_run("高亮 W20_HL ")
        r.font.highlight_color = WD_COLOR_INDEX.YELLOW
        r = p.add_run("红字 W20_RED")
        r.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        r.font.size = Pt(16)
        d.add_paragraph("正文段落 W20_BODY 中英 mixed")

    save("W20_styles_basic.docx", w20)

    def w21(d):
        d.add_heading("W21 多级列表 W21_LIST", 1)
        d.add_paragraph("一级 L1_A", style="List Bullet")
        d.add_paragraph("二级 L2_A", style="List Bullet 2")
        d.add_paragraph("二级 L2_B", style="List Bullet 2")
        d.add_paragraph("一级 L1_B", style="List Bullet")
        d.add_paragraph("步骤一 N1", style="List Number")
        d.add_paragraph("步骤二 N2", style="List Number")
        d.add_paragraph("收尾 W21_END")

    save("W21_multilevel_list.docx", w21)

    def w22(d):
        d.add_heading("W22 基础表加强 W22_TABLE", 1)
        d.add_paragraph("见表 W22_MARK")
        table(
            d,
            [
                ("区域 Region", "营收 Revenue", "同比 YoY", "备注 Note"),
                ("华北 North", "1,280", "+12%", "稳定"),
                ("华南 South", "960", "0%", "持平"),
                ("海外 Overseas", "430", "-5%", "承压"),
                ("合计 Total", "2,670", "+4%", "—"),
            ],
        )

    save("W22_table_rich.docx", w22)

    def w23(d):
        d.add_heading("W23 文+图加强 W23_IMG", 1)
        d.add_paragraph("开篇 W23_MARK")
        d.add_picture(str(assets["A"]), width=Inches(2))
        d.add_paragraph("图注 CAP_W23_A")
        d.add_picture(str(assets["B"]), width=Inches(1.8))
        d.add_paragraph("图注 CAP_W23_B")
        d.add_paragraph("结尾 W23_END")

    save("W23_images_basic.docx", w23)

    def w24(d):
        d.add_heading("W24 文图表齐 W24_MIX", 1)
        d.add_paragraph("综合 W24_MARK")
        d.add_picture(str(assets["C"]), width=Inches(2))
        d.add_paragraph("图注 CAP_W24_C")
        table(d, [("指标", "值", "阈值"), ("准确率", "96.2%", "↑"), ("召回率", "91.0%", "→")])
        d.add_paragraph("要点一", style="List Bullet")
        d.add_paragraph("要点二", style="List Bullet")
        d.add_paragraph("结尾 W24_END")

    save("W24_mix_basic.docx", w24)

    def w25(d):
        d.add_heading("W25 中英特殊字符 W25_CHARS", 1)
        d.add_paragraph("中文标点：【季度】「增长」《报告》——…… W25_CN")
        d.add_paragraph("EN: Q2's revenue = $1,280 (≈12%) — “North China”. W25_EN")
        d.add_paragraph("符号: ≤ ≥ ≠ ± × ÷ √ ∑ ∞ W25_SYM")
        d.add_paragraph("代码感: snake_case + camelCase + WORD_TOKEN W25_CODE")

    save("W25_special_chars.docx", w25)

    # ---- Wave3 深入 ----
    def w30(d):
        d.add_heading("W30 嵌套表外层 W30_NEST", 1)
        d.add_paragraph("外层说明 W30_MARK")
        outer = table(d, [["外层左 OUTER_L", "外层右 OUTER_R"], ["子表区", "备注区"]])
        # 在左下单元格内再插表（嵌套）
        cell = outer.rows[1].cells[0]
        cell.text = ""
        inner = cell.add_table(rows=3, cols=2)
        inner.style = "Table Grid"
        data = [("子项", "值"), ("华北", "1280"), ("INNER_MARK", "OK")]
        for i, row in enumerate(data):
            for j, v in enumerate(row):
                inner.rows[i].cells[j].text = v
        outer.rows[1].cells[1].text = "侧注 SIDE_NOTE"

    save("W30_nested_table.docx", w30)

    def w31(d):
        d.add_heading("W31 复杂合并 W31_MERGE", 1)
        t = table(
            d,
            [
                ["集团汇总 MERGE_TOP", "", "", ""],
                ["华北事业部", "收入", "1280", "+12%"],
                ["华北事业部", "成本", "800", "—"],
                ["海外事业部", "收入", "430", "-5%"],
                ["海外事业部", "成本", "300", "—"],
            ],
        )
        t.cell(0, 0).merge(t.cell(0, 3))
        t.cell(1, 0).merge(t.cell(2, 0))
        t.cell(3, 0).merge(t.cell(4, 0))
        d.add_paragraph("合并后文 W31_END")

    save("W31_complex_merge.docx", w31)

    def w32(d):
        d.add_heading("W32 多节分页 W32_SEC", 1)
        for i in range(1, 8):
            d.add_heading(f"第{i}章 CH{i:02d}", 1)
            d.add_heading(f"概述 {i}.1 SEC{i}_A", 2)
            d.add_paragraph(f"段落A W32_P{i}A 中英 content")
            d.add_heading(f"细节 {i}.2 SEC{i}_B", 2)
            d.add_paragraph(f"段落B W32_P{i}B")
            if i % 2 == 0:
                table(d, [("k", "v"), (f"r{i}", str(i * 10))])
            d.add_page_break()
        d.add_paragraph("全文收尾 W32_END")

    save("W32_multi_section_pagebreak.docx", w32)

    def w33(d):
        sec = d.sections[0]
        sec.header.paragraphs[0].text = "HDR_W33_公司机密"
        sec.footer.paragraphs[0].text = "FTR_W33_第页"
        d.add_heading("W33 页眉脚+正文 W33_HF", 1)
        d.add_paragraph("正文 BODY_W33")
        d.add_picture(str(assets["A"]), width=Inches(1.5))
        d.add_paragraph("图注 CAP_W33")
        table(d, [("区", "值"), ("华北", "100")])

    save("W33_header_footer_mix.docx", w33)

    def w34(d):
        d.add_heading("W34 超链接加强 W34_LINK", 1)
        p = d.add_paragraph()
        rid = p.part.relate_to(
            "https://example.com/w34?q=华北",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True,
        )
        hl = OxmlElement("w:hyperlink")
        hl.set(qn("r:id"), rid)
        nr = OxmlElement("w:r")
        te = OxmlElement("w:t")
        te.text = "点我 LINK_W34"
        nr.append(te)
        hl.append(nr)
        p._p.append(hl)
        d.add_paragraph("链接后文 W34_BODY 以及 email 感: user@example.com")

    save("W34_hyperlink_rich.docx", w34)

    def w35(d):
        d.add_heading("W35 压力包 W35_STRESS", 1)
        d.add_paragraph("压力开篇 W35_MARK")
        for i in range(1, 6):
            d.add_heading(f"块{i}", 2)
            d.add_picture(str(assets["A" if i % 2 else "B"]), width=Inches(1.4))
            d.add_paragraph(f"图注 CAP_W35_{i}")
            rows = [("列A", "列B", "列C")] + [(f"r{i}_{j}", str(i * j), f"m{j}") for j in range(1, 6)]
            table(d, rows)
        d.add_paragraph("压力结尾 W35_END")

    save("W35_stress_pack.docx", w35)

    def w36(d):
        d.add_paragraph("")  # almost empty

    save("W36_almost_empty.docx", w36)

    # 损坏/伪 docx
    bad = SAMPLES / "W37_corrupt.docx"
    bad.write_bytes(b"PK\x03\x04" + b"\x00" * 64 + b"NOT_A_REAL_DOCX W37_FAKE")
    names.append("W37_corrupt.docx")

    # 老 .doc 负面（假 OLE）
    fake_doc = SAMPLES / "W38_fake_old.doc"
    fake_doc.write_bytes(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"DOC_FAKE_W38" + b"\x00" * 200)
    names.append("W38_fake_old.doc")

    # 中文+空格文件名
    cn = SAMPLES / "W39 中文 空格 名.docx"
    doc = Document()
    doc.add_heading("W39 文件名边界 W39_NAME", 1)
    doc.add_paragraph("内容 W39_BODY")
    doc.save(cn)
    names.append(cn.name)

    return names


# ===================== PPT =====================
def gen_pptx(assets: dict[str, Path]) -> list[str]:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    names: list[str] = []

    def save(name, fn):
        p = SAMPLES / name
        prs = Presentation()
        fn(prs)
        prs.save(p)
        names.append(name)

    def title_body(prs, title, body, notes=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.text = body[0] if body else ""
        for line in body[1:]:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return slide

    def p20(prs):
        for i in range(1, 4):
            title_body(prs, f"P20 纯文{i} P20_T{i}", [f"正文 P20_TEXT_{i}", "中英 mixed"], notes=f"NOTE_P20_{i}")

    save("P20_text_basic.pptx", p20)

    def p21(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "P21 表格 P21_TT"
        table = slide.shapes.add_table(5, 4, Inches(0.5), Inches(1.8), Inches(9), Inches(2.5)).table
        data = [
            ("区域", "营收", "同比", "备注"),
            ("华北", "1280", "+12%", "稳"),
            ("华南", "960", "0%", "平"),
            ("海外", "430", "-5%", "压"),
            ("合计", "2670", "+4%", "—"),
        ]
        for r, row in enumerate(data):
            for c, v in enumerate(row):
                table.cell(r, c).text = v
        title_body(prs, "说明", ["P21_MARK"])

    save("P21_table_basic.pptx", p21)

    def p22(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "P22 多图 P22_IMG"
        slide.shapes.add_picture(str(assets["A"]), Inches(0.5), Inches(1.8), width=Inches(2.2))
        slide.shapes.add_picture(str(assets["B"]), Inches(3.2), Inches(1.8), width=Inches(2.2))
        slide.shapes.add_picture(str(assets["C"]), Inches(6.0), Inches(1.8), width=Inches(2.2))
        title_body(prs, "图注", ["CAP_P22_A CAP_P22_B CAP_P22_C", "P22_MARK"])

    save("P22_images_basic.pptx", p22)

    def p23(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "P23 综合 P23_MIX"
        slide.shapes.add_picture(str(assets["C"]), Inches(0.4), Inches(1.6), width=Inches(2.4))
        table = slide.shapes.add_table(3, 2, Inches(3.5), Inches(1.8), Inches(5), Inches(1.5)).table
        for r, row in enumerate([("指标", "值"), ("准确率", "96%"), ("召回率", "91%")]):
            for c, v in enumerate(row):
                table.cell(r, c).text = v
        title_body(prs, "说明", ["P23_MARK", "CAP_P23"])

    save("P23_mix_basic.pptx", p23)

    def p24(prs):
        slide = title_body(prs, "P24 样式列表 P24_STYLE", ["一级要点 P24_L1", "P24_MARK"])
        p = slide.shapes.title.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)
        p.font.size = Pt(32)
        body = slide.placeholders[1].text_frame
        p2 = body.add_paragraph()
        p2.text = "二级要点 P24_L2"
        p2.level = 1
        p2.font.color.rgb = RGBColor(0x05, 0x96, 0x69)
        p3 = body.add_paragraph()
        p3.text = "三级要点 P24_L3"
        p3.level = 2

    save("P24_styles_list.pptx", p24)

    def p25(prs):
        title_body(
            prs,
            "P25 中英字符 P25_CHARS",
            [
                "【季度】Revenue Q2 — “North China” +12%",
                "符号 ≤≥≠ 以及 snake_case TOKEN_P25",
                "P25_MARK",
            ],
        )

    save("P25_special_chars.pptx", p25)

    # Wave3
    def p30(prs):
        for i in range(1, 9):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"多节{i} P30_SEC{i}"
            if i % 2 == 0:
                table = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(5), Inches(1.4)).table
                for r, row in enumerate([("k", "v"), ("a", str(i)), ("b", str(i * 2))]):
                    for c, v in enumerate(row):
                        table.cell(r, c).text = v
            if i % 3 == 0:
                slide.shapes.add_picture(str(assets["A"]), Inches(6.5), Inches(2), width=Inches(2))
            slide.notes_slide.notes_text_frame.text = f"NOTE_P30_{i}"
        title_body(prs, "收尾", ["P30_END"])

    save("P30_multi_section.pptx", p30)

    def p31(prs):
        title_body(prs, "可见页 P31_VISIBLE", ["VISIBLE_P31"])
        hid = title_body(prs, "隐藏页 P31_HIDDEN", ["HIDDEN_P31"])
        hid._element.set("show", "0")
        title_body(prs, "再可见 P31_AFTER", ["AFTER_P31"])

    save("P31_hidden_slide.pptx", p31)

    def p32(prs):
        # 文本框自由形状（非占位符）
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tf = box.text_frame
        tf.text = "自由文本框 P32_FREEBOX"
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        box2 = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        box2.text_frame.text = "第二文本框 P32_BOX2\n换行内容 P32_MARK"

    save("P32_free_textbox.pptx", p32)

    def p33(prs):
        for i in range(1, 16):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"压力{i:02d} P33_S{i:02d}"
            if i % 2 == 0:
                table = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(5), Inches(1.3)).table
                for r, row in enumerate([("k", "v"), ("x", str(i)), ("y", str(i * 3))]):
                    for c, v in enumerate(row):
                        table.cell(r, c).text = v
            if i % 4 == 0:
                slide.shapes.add_picture(str(assets["B"]), Inches(6.5), Inches(2), width=Inches(2))
            slide.notes_slide.notes_text_frame.text = f"NOTE_P33_{i}"
        title_body(prs, "压力收尾", ["P33_END"])

    save("P33_stress_pack.pptx", p33)

    def p34(prs):
        prs.slides.add_slide(prs.slide_layouts[6])

    save("P34_blankish.pptx", p34)

    bad = SAMPLES / "P35_corrupt.pptx"
    bad.write_bytes(b"PK\x03\x04" + b"\x00" * 40 + b"FAKE_PPTX_P35")
    names.append("P35_corrupt.pptx")

    fake_ppt = SAMPLES / "P36_fake_old.ppt"
    fake_ppt.write_bytes(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"PPT_FAKE_P36" + b"\x00" * 100)
    names.append("P36_fake_old.ppt")

    cn = SAMPLES / "P37 演示 中文名.pptx"
    prs = Presentation()
    title_body(prs, "P37 文件名 P37_NAME", ["内容 P37_BODY"])
    prs.save(cn)
    names.append(cn.name)

    return names


# ===================== Excel =====================
def gen_excel(assets: dict[str, Path]) -> list[str]:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.comments import Comment
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Font, PatternFill, numbers
    from openpyxl.utils import get_column_letter
    from openpyxl.workbook.protection import WorkbookProtection
    import xlwt

    names: list[str] = []

    def save_xlsx(name, fn):
        p = SAMPLES / name
        wb = Workbook()
        fn(wb)
        wb.save(p)
        names.append(name)

    def e20(wb):
        ws = wb.active
        ws.title = "说明"
        ws["A1"] = "E20 纯文 E20_TEXT"
        ws["A2"] = "华北增长 E20_CN"
        ws["A3"] = "Follow-up E20_EN"

    save_xlsx("E20_text_basic.xlsx", e20)

    def e21(wb):
        ws = wb.active
        ws.title = "收入"
        ws["A1"] = "E21 表 E21_TT"
        headers = ["区域", "营收", "同比", "备注"]
        for c, h in enumerate(headers, 1):
            cell = ws.cell(3, c, h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="2563EB")
        for i, row in enumerate(
            [("华北", 1280, "+12%", "稳"), ("华南", 960, "0%", "平"), ("海外", 430, "-5%", "压")], 4
        ):
            for j, v in enumerate(row, 1):
                ws.cell(i, j, v)
        ws["A7"] = "E21_MARK"

    save_xlsx("E21_table_basic.xlsx", e21)

    def e22(wb):
        ws = wb.active
        ws.title = "图"
        ws["A1"] = "E22 图 E22_IMG"
        ws["A2"] = "E22_MARK"
        for key, cell in [("A", "A4"), ("B", "D4")]:
            img = XLImage(str(assets[key]))
            img.width, img.height = 140, 48
            ws.add_image(img, cell)
        ws["A10"] = "CAP_E22"

    save_xlsx("E22_images_basic.xlsx", e22)

    def e23(wb):
        ws = wb.active
        ws.title = "综合"
        ws["A1"] = "E23 综合 E23_MIX"
        ws["A2"] = "E23_MARK"
        img = XLImage(str(assets["C"]))
        img.width, img.height = 160, 65
        ws.add_image(img, "A3")
        ws["A10"] = "指标"
        ws["B10"] = "值"
        ws["A11"] = "准确率"
        ws["B11"] = "96%"
        ws["A12"] = "CAP_E23"

    save_xlsx("E23_mix_basic.xlsx", e23)

    def e24(wb):
        ws = wb.active
        ws.title = "混排"
        ws["A1"] = "E24 Q2 Revenue 季度【华北】 E24_CHARS"
        ws["A2"] = "North China（华北）= 1,280 ≈ +12%"
        ws["A3"] = "符号 ≤≥≠ snake_case E24_MARK"

    save_xlsx("E24_special_chars.xlsx", e24)

    def e25(wb):
        ws = wb.active
        ws.title = "并排"
        ws["A1"] = "左表 E25_LEFT"
        ws["A2"] = "区"
        ws["B2"] = "值"
        ws["A3"] = "华北"
        ws["B3"] = 1280
        ws["D1"] = "右表 E25_RIGHT"
        ws["D2"] = "项"
        ws["E2"] = "额"
        ws["D3"] = "人力"
        ws["E3"] = 200

    save_xlsx("E25_side_tables.xlsx", e25)

    # Wave3
    def e30(wb):
        ws = wb.active
        ws.title = "合并"
        ws.merge_cells("A1:F1")
        ws["A1"] = "E30 复杂合并 E30_MERGE_HDR"
        ws.merge_cells("A2:A5")
        ws["A2"] = "华北组"
        ws["B2"] = "收入"
        ws["C2"] = 1280
        ws["D2"] = "+12%"
        ws["E2"] = "备注A"
        ws["F2"] = "OK"
        ws["B3"] = "成本"
        ws["C3"] = 800
        ws["D3"] = "—"
        ws["E3"] = "备注B"
        ws["F3"] = "OK"
        ws["B4"] = "利润"
        ws["C4"] = 480
        ws["D4"] = "—"
        ws["E4"] = "备注C"
        ws["F4"] = "OK"
        ws["B5"] = "人数"
        ws["C5"] = 32
        ws["D5"] = "—"
        ws["E5"] = "备注D"
        ws["F5"] = "OK"
        ws.merge_cells("A6:A7")
        ws["A6"] = "海外组"
        ws["B6"] = "收入"
        ws["C6"] = 430
        ws["D6"] = "-5%"
        ws["E6"] = "承压"
        ws["F6"] = "WARN"
        ws["B7"] = "成本"
        ws["C7"] = 300
        ws["D7"] = "—"
        ws["E7"] = "控本"
        ws["F7"] = "OK"
        ws["A8"] = "E30_END"

    save_xlsx("E30_complex_merge.xlsx", e30)

    def e31(wb):
        names_sheets = ["华北", "华南", "海外", "汇总", "HiddenSrc"]
        for i, name in enumerate(names_sheets):
            ws = wb.active if i == 0 else wb.create_sheet(name)
            if i == 0:
                ws.title = name
            ws["A1"] = f"SHEET_{name}"
            ws["A2"] = "营收"
            ws["B2"] = 100 * (i + 1)
            ws["A3"] = f"MARK_{name}"
        wb["HiddenSrc"].sheet_state = "hidden"
        wb["汇总"]["A4"] = "E31_MULTI"

    save_xlsx("E31_multi_sheet_hidden.xlsx", e31)

    def e32(wb):
        ws = wb.active
        ws.title = "类型"
        ws["A1"] = "E32 类型 E32_TYPES"
        ws["A2"] = "日期"
        ws["B2"] = "2024-06-15"
        ws["B2"].number_format = "YYYY-MM-DD"
        ws["A3"] = "百分数"
        ws["B3"] = 0.125
        ws["B3"].number_format = "0.0%"
        ws["A4"] = "货币"
        ws["B4"] = 1280.5
        ws["B4"].number_format = '"¥"#,##0.00'
        ws["A5"] = "布尔"
        ws["B5"] = True
        ws["A6"] = "公式"
        ws["B6"] = "=1/0"
        ws["A7"] = "空公式"
        ws["B7"] = "=IF(1=1,\"E32_OK\",\"no\")"
        ws["A8"] = "E32_MARK"

    save_xlsx("E32_types_formula.xlsx", e32)

    def e33(wb):
        ws = wb.active
        ws.title = "宽高"
        ws["A1"] = "E33 宽表高表 E33_WIDE"
        for c in range(1, 31):
            ws.cell(2, c, f"列{c}")
            ws.cell(3, c, c * 7)
        for r in range(4, 25):
            ws.cell(r, 1, f"行{r}")
            ws.cell(r, 2, r * 3)
        ws["A26"] = "E33_MARK"

    save_xlsx("E33_wide_tall.xlsx", e33)

    def e34(wb):
        ws = wb.active
        ws.title = "批注图"
        ws["A1"] = "E34 批注 E34_COMMENT"
        ws["A2"] = "华北"
        ws["A2"].comment = Comment("批注内容 COMMENT_E34", "tester")
        img = XLImage(str(assets["A"]))
        img.width, img.height = 120, 40
        ws.add_image(img, "C2")
        ws["A5"] = "E34_MARK"
        # 简单图表（抽取时可能只剩数据）
        ws["A7"] = "区"
        ws["B7"] = "值"
        ws["A8"] = "华北"
        ws["B8"] = 10
        ws["A9"] = "海外"
        ws["B9"] = 5
        chart = BarChart()
        chart.title = "CHART_E34"
        data = Reference(ws, min_col=2, min_row=7, max_row=9)
        cats = Reference(ws, min_col=1, min_row=8, max_row=9)
        chart.add_data(data, titles_from_data=True)
        chart.set_categories(cats)
        ws.add_chart(chart, "D7")

    save_xlsx("E34_comment_chart_image.xlsx", e34)

    def e35(wb):
        for si in range(1, 5):
            ws = wb.active if si == 1 else wb.create_sheet(f"S{si}")
            if si == 1:
                ws.title = "S1"
            ws["A1"] = f"E35 压力 S{si} E35_S{si}"
            ws["A2"] = "区"
            ws["B2"] = "值"
            for r in range(3, 12):
                ws.cell(r, 1, f"R{si}_{r}")
                ws.cell(r, 2, si * r)
            ws["D1"] = f"右表 E35_R{si}"
            ws["D2"] = "项"
            ws["E2"] = "额"
            ws["D3"] = "云"
            ws["E3"] = 50 * si
            try:
                img = XLImage(str(assets["A"]))
                img.width, img.height = 100, 35
                ws.add_image(img, "A14")
            except Exception:
                pass
        wb.create_sheet("END")["A1"] = "E35_END"

    save_xlsx("E35_stress_pack.xlsx", e35)

    def e36(wb):
        wb.active["A1"] = ""

    save_xlsx("E36_emptyish.xlsx", e36)

    # xls 多 sheet
    p = SAMPLES / "E37_xls_multi.xls"
    book = xlwt.Workbook()
    for name, mark, val in [("华北", "XLS_N", 1280), ("海外", "XLS_O", 430)]:
        sh = book.add_sheet(name)
        sh.write(0, 0, f"E37 {name} {mark}")
        sh.write(1, 0, "区域")
        sh.write(1, 1, "营收")
        sh.write(2, 0, name)
        sh.write(2, 1, val)
        sh.write(3, 0, "E37_MARK")
    book.save(str(p))
    names.append("E37_xls_multi.xls")

    # 损坏 xlsx
    bad = SAMPLES / "E38_corrupt.xlsx"
    bad.write_bytes(b"PK\x03\x04" + b"\x00" * 30 + b"FAKE_XLSX_E38")
    names.append("E38_corrupt.xlsx")

    # 中文文件名
    cn = SAMPLES / "E39 评分表 中文.xlsx"
    wb = Workbook()
    wb.active["A1"] = "E39 文件名 E39_NAME"
    wb.active["A2"] = "内容 E39_BODY"
    wb.save(cn)
    names.append(cn.name)

    # 伪加密提示：普通文件但标记（真加密难造，用 workbook protection 元数据）
    def e40(wb):
        ws = wb.active
        ws.title = "保护"
        ws["A1"] = "E40 工作簿保护标记 E40_PROT"
        ws["A2"] = "E40_MARK"
        wb.security = WorkbookProtection(lockStructure=True, workbookPassword="test")

    save_xlsx("E40_workbook_protect_meta.xlsx", e40)

    return names


# ===================== MSG =====================
def gen_msg() -> list[str]:
    names: list[str] = []
    SAMPLES.mkdir(parents=True, exist_ok=True)

    src = DEEP_MSG if DEEP_MSG.exists() else PREV_MSG
    if src.exists():
        dest = SAMPLES / "M20_official.msg"
        shutil.copy2(src, dest)
        names.append(dest.name)
        dest2 = SAMPLES / "M21 邮件 中文名.msg"
        shutil.copy2(src, dest2)
        names.append(dest2.name)

    # 截断损坏
    if src.exists():
        data = src.read_bytes()
        trunc = SAMPLES / "M30_truncated.msg"
        trunc.write_bytes(data[: max(64, len(data) // 10)])
        names.append(trunc.name)

    # 伪 OLE
    fake = SAMPLES / "M31_fake_ole.msg"
    fake.write_bytes(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1" + b"\x00" * 200)
    names.append(fake.name)

    # 完全随机
    rnd = SAMPLES / "M32_random.bin.msg"
    rnd.write_bytes(b"NOT_OLE_M32_" + bytes(range(64)))
    names.append(rnd.name)

    # 错误后缀：真实 msg 改名为 .docx（负面/探测）
    if src.exists():
        mismatch = SAMPLES / "M33_msg_as.docx"
        shutil.copy2(src, mismatch)
        names.append(mismatch.name)

    return names


def main():
    assets = make_assets()
    w = gen_word(assets)
    p = gen_pptx(assets)
    e = gen_excel(assets)
    m = gen_msg()
    print(f"word={len(w)} ppt={len(p)} excel={len(e)} msg={len(m)} total={len(w)+len(p)+len(e)+len(m)}")
    for lst, tag in [(w, "W"), (p, "P"), (e, "E"), (m, "M")]:
        print(tag, ",".join(lst))


if __name__ == "__main__":
    main()
