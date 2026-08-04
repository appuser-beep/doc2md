"""Office Wave4 再加深 + PDF 穷举样例生成。"""

from __future__ import annotations

import io
import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

ROOT = Path(__file__).resolve().parent
SAMPLES = ROOT / "samples"
ASSETS = ROOT / "assets"
MSG_SRC = Path(__file__).resolve().parents[1] / "office_deep" / "samples" / "MSG01_official.msg"


def _font_img(size: int = 16):
    for name in (r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\arial.ttf"):
        if Path(name).exists():
            try:
                return ImageFont.truetype(name, size)
            except OSError:
                pass
    return ImageFont.load_default()


def make_assets() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    out = {}
    for key, color, size in [
        ("A", (37, 99, 235), (160, 55)),
        ("B", (5, 150, 105), (160, 55)),
        ("C", (220, 38, 38), (200, 80)),
        ("D", (100, 100, 100), (120, 50)),
        ("SCAN", (240, 240, 240), (400, 200)),
    ]:
        p = ASSETS / f"w4_{key}.png"
        im = Image.new("RGB", size, color)
        dr = ImageDraw.Draw(im)
        label = "SCANNED_PDF_MARKER" if key == "SCAN" else f"IMG_{key}"
        dr.text((10, size[1] // 2 - 8), label, fill=(20, 20, 20) if key == "SCAN" else (255, 255, 255), font=_font_img(14))
        im.save(p)
        out[key] = p
    return out


# ===================== Office Wave4 =====================
def gen_office_wave4(assets: dict[str, Path]) -> list[str]:
    from docx import Document
    from docx.enum.text import WD_COLOR_INDEX
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
    from docx.shared import Inches, Pt, RGBColor, Twips
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    import xlwt

    names: list[str] = []
    SAMPLES.mkdir(parents=True, exist_ok=True)

    def save_docx(name, fn):
        p = SAMPLES / name
        doc = Document()
        fn(doc)
        doc.save(p)
        names.append(name)

    def tbl(doc, rows):
        t = doc.add_table(rows=len(rows), cols=len(rows[0]))
        t.style = "Table Grid"
        for i, row in enumerate(rows):
            for j, v in enumerate(row):
                t.rows[i].cells[j].text = str(v)
        return t

    # --- Word Wave4 ---
    def w40(d):
        d.add_heading("W40 脚注感标注 W40_NOTE", 1)
        d.add_paragraph("正文含脚注引用感 W40_BODY¹")
        d.add_paragraph("注1：脚注内容 FOOTNOTE_W40")
        d.add_paragraph("批注感：COMMENT_W40 已确认")
        p = d.add_paragraph()
        r = p.add_run("删除线 DEL_W40 ")
        r.font.strike = True
        r = p.add_run("保留 KEEP_W40")

    save_docx("W40_footnote_strike.docx", w40)

    def w41(d):
        d.add_heading("W41 文本框替代：引用块 W41_QUOTE", 1)
        d.add_paragraph("前置 W41_MARK")
        d.add_paragraph("「引用内容 QUOTE_W41：华北增长12%」")
        d.add_paragraph("后置 W41_END")

    save_docx("W41_quote_block.docx", w41)

    def w42(d):
        # 分节：不同页眉
        d.add_heading("W42 多节页眉 W42_SEC", 1)
        d.sections[0].header.paragraphs[0].text = "HDR_SEC1_W42"
        d.add_paragraph("第一节 BODY_SEC1_W42")
        d.add_page_break()
        # new section
        new_sec = d.add_section()
        new_sec.header.is_linked_to_previous = False
        new_sec.header.paragraphs[0].text = "HDR_SEC2_W42"
        d.add_heading("第二节", 1)
        d.add_paragraph("第二节 BODY_SEC2_W42")
        tbl(d, [("区", "值"), ("华北", "100")])

    save_docx("W42_multi_section_headers.docx", w42)

    def w43(d):
        d.add_heading("W43 并排感双表 W43_DUAL", 1)
        d.add_paragraph("表A说明 W43_A")
        tbl(d, [("A区", "A值"), ("华北", "1280"), ("A_END", "1")])
        d.add_paragraph("表B说明 W43_B")
        tbl(d, [("B项", "B额"), ("人力", "200"), ("B_END", "2")])
        d.add_picture(str(assets["A"]), width=Inches(1.5))
        d.add_paragraph("图注 CAP_W43")

    save_docx("W43_dual_table_image.docx", w43)

    def w44(d):
        d.add_heading("W44 超长段落压力 W44_LONG", 1)
        chunk = "华北增长。Overseas decline. " * 40
        d.add_paragraph("开篇 W44_MARK " + chunk)
        d.add_paragraph("中段 W44_MID " + chunk)
        d.add_paragraph("收尾 W44_END")

    save_docx("W44_long_paragraphs.docx", w44)

    def w45(d):
        d.add_heading("W45 emoji与零宽 W45_EMOJI", 1)
        d.add_paragraph("状态 ✅ 警告 ⚠️ 完成 🎉 W45_MARK")
        # zero-width space between
        d.add_paragraph("零宽\u200b连接 W45_ZWSP")
        d.add_paragraph("全角数字：１２８０ 与半角 1280 W45_NUM")

    save_docx("W45_emoji_zwsp.docx", w45)

    # --- PPT Wave4 ---
    def save_pptx(name, fn):
        p = SAMPLES / name
        prs = Presentation()
        fn(prs)
        prs.save(p)
        names.append(name)

    def tb(prs, title, body, notes=None):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.text = body[0] if body else ""
        for line in body[1:]:
            p = tf.add_paragraph()
            p.text = line
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return slide

    def p40(prs):
        # 两栏感：两个文本框
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        left = slide.shapes.add_textbox(PInches(0.5), PInches(1), PInches(4), PInches(4))
        left.text_frame.text = "左栏 LEFT_P40\n华北 1280\nP40_MARK"
        right = slide.shapes.add_textbox(PInches(5), PInches(1), PInches(4), PInches(4))
        right.text_frame.text = "右栏 RIGHT_P40\n海外 430"
        tb(prs, "说明页", ["P40_END"], notes="NOTE_P40")

    save_pptx("P40_two_column_boxes.pptx", p40)

    def p41(prs):
        for i in range(1, 6):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"P41 混合{i} P41_S{i}"
            slide.shapes.add_picture(str(assets["A"]), PInches(0.5), PInches(1.8), width=PInches(2))
            table = slide.shapes.add_table(3, 2, PInches(3.5), PInches(1.8), PInches(5), PInches(1.4)).table
            for r, row in enumerate([("k", "v"), ("华北", str(100 * i)), ("海外", str(50 * i))]):
                for c, v in enumerate(row):
                    table.cell(r, c).text = v
            slide.notes_slide.notes_text_frame.text = f"NOTE_P41_{i}"
        tb(prs, "收尾", ["P41_END"])

    save_pptx("P41_multi_mix_stress.pptx", p41)

    def p42(prs):
        slide = tb(prs, "P42 超长备注 P42_NOTES", ["正文 P42_BODY"])
        slide.notes_slide.notes_text_frame.text = "NOTE_P42_LONG " + ("备注扩展。 " * 50) + " NOTE_P42_END"

    save_pptx("P42_long_notes.pptx", p42)

    def p43(prs):
        tb(prs, "P43 emoji P43_EMOJI", ["状态 ✅ 完成 P43_MARK", "警告 ⚠️ Overseas"])

    save_pptx("P43_emoji.pptx", p43)

    # --- Excel Wave4 ---
    def save_xlsx(name, fn):
        p = SAMPLES / name
        wb = Workbook()
        fn(wb)
        wb.save(p)
        names.append(name)

    def e40(wb):
        ws = wb.active
        ws.title = "并排合并"
        ws.merge_cells("A1:B1")
        ws["A1"] = "左块 E40_LEFT"
        ws["A2"] = "区"
        ws["B2"] = "值"
        ws["A3"] = "华北"
        ws["B3"] = 1280
        ws.merge_cells("D1:E1")
        ws["D1"] = "右块 E40_RIGHT"
        ws["D2"] = "项"
        ws["E2"] = "额"
        ws["D3"] = "人力"
        ws["E3"] = 200
        ws["A5"] = "E40_MARK"

    save_xlsx("E40_side_merge.xlsx", e40)

    def e41(wb):
        # 很多 sheet
        for i in range(1, 13):
            ws = wb.active if i == 1 else wb.create_sheet(f"S{i:02d}")
            if i == 1:
                ws.title = "S01"
            ws["A1"] = f"E41_SHEET_{i:02d}"
            ws["A2"] = "值"
            ws["B2"] = i * 11
        wb.create_sheet("END")["A1"] = "E41_END"

    save_xlsx("E41_many_sheets.xlsx", e41)

    def e42(wb):
        ws = wb.active
        ws.title = "稀疏宽"
        ws["A1"] = "E42 稀疏宽表 E42_SPARSE"
        for c in range(1, 16):
            ws.cell(2, c, f"H{c}")
        ws["A3"] = "华北"
        ws["B3"] = 1280
        ws["O3"] = "尾列 E42_TAIL"  # col 15
        ws["A4"] = "E42_MARK"

    save_xlsx("E42_sparse_wide.xlsx", e42)

    def e43(wb):
        ws = wb.active
        ws.title = "样式边框"
        thin = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )
        ws["A1"] = "E43 样式 E43_STYLE"
        for c, h in enumerate(["区域", "营收"], 1):
            cell = ws.cell(2, c, h)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill("solid", fgColor="DC2626")
            cell.border = thin
            cell.alignment = Alignment(horizontal="center")
        ws["A3"] = "华北"
        ws["B3"] = 1280
        ws["A3"].border = thin
        ws["B3"].border = thin
        ws["A4"] = "E43_MARK"

    save_xlsx("E43_styled_table.xlsx", e43)

    def e44(wb):
        ws = wb.active
        ws.title = "图文"
        ws["A1"] = "E44 多图 E44_IMG"
        for i, key in enumerate(["A", "B", "C"]):
            img = XLImage(str(assets[key]))
            img.width, img.height = 120, 40
            ws.add_image(img, f"A{3 + i * 4}")
        ws["A16"] = "E44_MARK"

    save_xlsx("E44_multi_image.xlsx", e44)

    # xls 宽一点
    p = SAMPLES / "E45_xls_rich.xls"
    book = xlwt.Workbook()
    sh = book.add_sheet("收入")
    sh.write(0, 0, "E45 XLS富表 E45_XLS")
    for c, h in enumerate(["区域", "营收", "同比"]):
        sh.write(1, c, h)
    for r, row in enumerate([("华北", 1280, "+12%"), ("海外", 430, "-5%")]):
        for c, v in enumerate(row):
            sh.write(2 + r, c, v)
    sh.write(4, 0, "E45_MARK")
    sh2 = book.add_sheet("备注")
    sh2.write(0, 0, "E45_NOTE_SHEET")
    book.save(str(p))
    names.append("E45_xls_rich.xls")

    # MSG copies with stress names
    if MSG_SRC.exists():
        dest = SAMPLES / "M40_official_wave4.msg"
        shutil.copy2(MSG_SRC, dest)
        names.append(dest.name)
        dest2 = SAMPLES / "M41（括号）邮件.msg"
        shutil.copy2(MSG_SRC, dest2)
        names.append(dest2.name)

    # 损坏 msg / 空文件
    (SAMPLES / "M42_empty.msg").write_bytes(b"")
    names.append("M42_empty.msg")

    # Office 中文路径综合
    cn = SAMPLES / "W46（综合）中文名.docx"
    doc = Document()
    doc.add_heading("W46 中文名综合 W46_NAME", 1)
    doc.add_paragraph("内容 W46_BODY")
    tbl(doc, [("k", "v"), ("华北", "1")])
    doc.add_picture(str(assets["A"]), width=Inches(1.2))
    doc.add_paragraph("图注 CAP_W46")
    doc.save(cn)
    names.append(cn.name)

    return names


# ===================== PDF Exhaust =====================
def _register_font():
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for name, path in [("MSYH", r"C:\Windows\Fonts\msyh.ttc"), ("SIMSUN", r"C:\Windows\Fonts\simsun.ttc")]:
        if Path(path).exists():
            try:
                pdfmetrics.registerFont(TTFont(name, path))
                return name
            except Exception:
                continue
    return "Helvetica"


def gen_pdf(assets: dict[str, Path]) -> list[str]:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import (
        Image as RLImage,
        PageBreak,
        Paragraph,
        SimpleDocTemplate,
        Spacer,
        Table,
        TableStyle,
    )
    from reportlab.pdfgen import canvas as pdfcanvas

    names: list[str] = []
    font = _register_font()
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("T", parent=styles["Heading1"], fontName=font, fontSize=16, leading=22)
    h2_style = ParagraphStyle("H2", parent=styles["Heading2"], fontName=font, fontSize=13, leading=18)
    body_style = ParagraphStyle("B", parent=styles["Normal"], fontName=font, fontSize=11, leading=16)

    def save_flow(name, build_story):
        p = SAMPLES / name
        doc = SimpleDocTemplate(str(p), pagesize=A4, topMargin=18 * mm, bottomMargin=18 * mm)
        story = []
        build_story(story)
        doc.build(story)
        names.append(name)

    def pdf01(story):
        story.append(Paragraph("PDF01 纯文字 PDF_TEXT", title_style))
        story.append(Paragraph("华北增长12%。海外-5%。PDF_BODY_CN", body_style))
        story.append(Paragraph("Action item follow-up. PDF_BODY_EN", body_style))
        story.append(Paragraph("要点：稳定扩产 PDF_END", body_style))

    save_flow("PDF01_text_only.pdf", pdf01)

    def pdf02(story):
        story.append(Paragraph("PDF02 文+表 PDF_TT", title_style))
        story.append(Paragraph("见表 PDF_TT_MARK", body_style))
        data = [
            ["区域", "营收", "同比"],
            ["华北", "1280", "+12%"],
            ["华南", "960", "0%"],
            ["海外", "430", "-5%"],
        ]
        t = Table(data, colWidths=[40 * mm, 30 * mm, 30 * mm])
        t.setStyle(
            TableStyle(
                [
                    ("FONTNAME", (0, 0), (-1, -1), font),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563EB")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                ]
            )
        )
        story.append(t)

    save_flow("PDF02_text_table.pdf", pdf02)

    def pdf03(story):
        story.append(Paragraph("PDF03 文+图 PDF_TI", title_style))
        story.append(Paragraph("见图 PDF_TI_MARK", body_style))
        story.append(RLImage(str(assets["A"]), width=60 * mm, height=20 * mm))
        story.append(Paragraph("图注 CAP_PDF_A", body_style))

    save_flow("PDF03_text_image.pdf", pdf03)

    def pdf04(story):
        story.append(Paragraph("PDF04 文图表 PDF_MIX", title_style))
        story.append(Paragraph("综合 PDF_MIX_MARK", body_style))
        story.append(RLImage(str(assets["C"]), width=55 * mm, height=22 * mm))
        story.append(Paragraph("图注 CAP_PDF_C", body_style))
        data = [["指标", "值"], ["准确率", "96%"], ["召回率", "91%"]]
        t = Table(data, colWidths=[40 * mm, 30 * mm])
        t.setStyle(TableStyle([("FONTNAME", (0, 0), (-1, -1), font), ("GRID", (0, 0), (-1, -1), 0.5, colors.grey)]))
        story.append(t)
        story.append(Paragraph("结尾 PDF_MIX_END", body_style))

    save_flow("PDF04_text_image_table.pdf", pdf04)

    def pdf05(story):
        story.append(Paragraph("PDF05 多表 PDF_MT", title_style))
        for i, mark in enumerate(["PDF_MT1", "PDF_MT2", "PDF_MT3"], 1):
            story.append(Paragraph(f"表{i} {mark}", h2_style))
            data = [["k", "v"], [f"a{i}", str(i * 10)], [f"b{i}", str(i * 20)]]
            t = Table(data, colWidths=[40 * mm, 30 * mm])
            t.setStyle(TableStyle([("FONTNAME", (0, 0), (-1, -1), font), ("GRID", (0, 0), (-1, -1), 0.5, colors.grey)]))
            story.append(t)
            story.append(Spacer(1, 8))

    save_flow("PDF05_multi_table.pdf", pdf05)

    def pdf06(story):
        story.append(Paragraph("PDF06 多图 PDF_MI", title_style))
        story.append(Paragraph("多图 PDF_MI_MARK", body_style))
        for key, cap in [("A", "CAP_PDF_A1"), ("B", "CAP_PDF_B1"), ("D", "CAP_PDF_D1")]:
            story.append(RLImage(str(assets[key]), width=50 * mm, height=18 * mm))
            story.append(Paragraph(f"图注 {cap}", body_style))

    save_flow("PDF06_multi_image.pdf", pdf06)

    def pdf07(story):
        story.append(Paragraph("PDF07 多页多节 PDF_SEC", title_style))
        for i in range(1, 6):
            story.append(Paragraph(f"第{i}章 CH{i:02d}", h2_style))
            story.append(Paragraph(f"内容 PDF_SEC_{i} 中英 content", body_style))
            if i < 5:
                story.append(PageBreak())
        story.append(Paragraph("收尾 PDF_SEC_END", body_style))

    save_flow("PDF07_multipage_sections.pdf", pdf07)

    def pdf08(story):
        story.append(Paragraph("PDF08 中英混排 PDF_LANG", title_style))
        story.append(
            Paragraph(
                "Q2 Revenue 季度营收：North China（华北）+12%。Overseas（海外）-5%。PDF_LANG_MARK",
                body_style,
            )
        )
        story.append(Paragraph("符号 ≤ ≥ ≠ ± —— 「引号」 PDF_LANG_SYM", body_style))

    save_flow("PDF08_mixed_lang.pdf", pdf08)

    def pdf09(story):
        # 两栏：用宽表模拟左右栏文本不稳，改用 canvas 更准；此处用双列表
        story.append(Paragraph("PDF09 双栏感 PDF_COL", title_style))
        data = [
            [Paragraph("左栏 LEFT_PDF09<br/>华北增长 PDF_COL_L", body_style),
             Paragraph("右栏 RIGHT_PDF09<br/>海外承压 PDF_COL_R", body_style)]
        ]
        t = Table(data, colWidths=[80 * mm, 80 * mm])
        t.setStyle(TableStyle([("VALIGN", (0, 0), (-1, -1), "TOP"), ("BOX", (0, 0), (-1, -1), 0.5, colors.grey)]))
        story.append(t)
        story.append(Paragraph("PDF_COL_MARK", body_style))

    save_flow("PDF09_two_column.pdf", pdf09)

    def pdf10(story):
        story.append(Paragraph("PDF10 压力包 PDF_STRESS", title_style))
        story.append(Paragraph("开篇 PDF_STRESS_MARK", body_style))
        for i in range(1, 5):
            story.append(Paragraph(f"块{i}", h2_style))
            story.append(RLImage(str(assets["A" if i % 2 else "B"]), width=45 * mm, height=16 * mm))
            story.append(Paragraph(f"图注 CAP_PDF_ST_{i}", body_style))
            data = [["k", "v"], [f"r{i}a", str(100 * i)], [f"r{i}b", str(200 * i)]]
            t = Table(data, colWidths=[35 * mm, 25 * mm])
            t.setStyle(TableStyle([("FONTNAME", (0, 0), (-1, -1), font), ("GRID", (0, 0), (-1, -1), 0.5, colors.grey)]))
            story.append(t)
            if i < 4:
                story.append(PageBreak())
        story.append(Paragraph("压力结尾 PDF_STRESS_END", body_style))

    save_flow("PDF10_stress_pack.pdf", pdf10)

    def pdf11(story):
        story.append(Paragraph("PDF11 宽表 PDF_WIDE", title_style))
        header = [f"列{i}" for i in range(1, 9)]
        row = [str(i * 11) for i in range(1, 9)]
        data = [header, row, ["华北"] + ["x"] * 7]
        t = Table(data, colWidths=[22 * mm] * 8)
        t.setStyle(
            TableStyle(
                [
                    ("FONTNAME", (0, 0), (-1, -1), font),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.4, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E5E7EB")),
                ]
            )
        )
        story.append(t)
        story.append(Paragraph("PDF_WIDE_MARK", body_style))

    save_flow("PDF11_wide_table.pdf", pdf11)

    def pdf12(story):
        story.append(Paragraph("PDF12 合并感表 PDF_MERGE", title_style))
        # reportlab span
        data = [
            ["合计区 MERGE_HDR", "", ""],
            ["华北", "1280", "+12%"],
            ["海外", "430", "-5%"],
        ]
        t = Table(data, colWidths=[40 * mm, 30 * mm, 30 * mm])
        t.setStyle(
            TableStyle(
                [
                    ("FONTNAME", (0, 0), (-1, -1), font),
                    ("SPAN", (0, 0), (2, 0)),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("BACKGROUND", (0, 0), (2, 0), colors.HexColor("#FEF3C7")),
                ]
            )
        )
        story.append(t)
        story.append(Paragraph("PDF_MERGE_MARK", body_style))

    save_flow("PDF12_merged_table.pdf", pdf12)

    # 几乎空
    def pdf13(story):
        story.append(Spacer(1, 10))

    save_flow("PDF13_almost_empty.pdf", pdf13)

    # 扫描件：整页图片，无文字层
    p = SAMPLES / "PDF14_scanned_like.pdf"
    c = pdfcanvas.Canvas(str(p), pagesize=A4)
    c.drawImage(str(assets["SCAN"]), 50, 500, width=400, height=200, mask="auto")
    c.showPage()
    c.save()
    names.append("PDF14_scanned_like.pdf")

    # 损坏
    (SAMPLES / "PDF15_corrupt.pdf").write_bytes(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n1 0 obj<<>>endobj\ntrailer<<>>\n")
    names.append("PDF15_corrupt.pdf")

    # 伪 PDF
    (SAMPLES / "PDF16_fake.pdf").write_bytes(b"NOT_A_PDF PDF16_FAKE" + b"\x00" * 100)
    names.append("PDF16_fake.pdf")

    # 中文文件名
    cn = SAMPLES / "PDF17 报告 中文名.pdf"
    doc = SimpleDocTemplate(str(cn), pagesize=A4)
    doc.build(
        [
            Paragraph("PDF17 中文文件名 PDF17_NAME", title_style),
            Paragraph("内容 PDF17_BODY", body_style),
        ]
    )
    names.append(cn.name)

    # 加密 PDF（reportlab 可用 encrypt）
    try:
        from reportlab.lib.utils import open_for_read
        from reportlab.pdfgen.canvas import Canvas

        enc = SAMPLES / "PDF18_encrypted.pdf"
        c = Canvas(str(enc), pagesize=A4)
        # encrypt API
        c.encrypt(userPassword="user", ownerPassword="owner", canPrint=0)
        c.setFont(font, 14)
        c.drawString(72, 720, "PDF18 加密 PDF18_ENC")
        c.drawString(72, 700, "不应明文抽出 PDF18_SECRET")
        c.save()
        names.append("PDF18_encrypted.pdf")
    except Exception as e:
        print("encrypt skip", e)

    # 超多页
    def pdf19(story):
        story.append(Paragraph("PDF19 多页压力 PDF19_MANY", title_style))
        for i in range(1, 13):
            story.append(Paragraph(f"页内容 PDF19_P{i:02d}", body_style))
            if i < 12:
                story.append(PageBreak())
        story.append(Paragraph("PDF19_END", body_style))

    save_flow("PDF19_many_pages.pdf", pdf19)

    # 仅表无说明
    def pdf20(story):
        data = [["区域", "值"], ["华北", "100"], ["PDF20_MARK", "1"]]
        t = Table(data, colWidths=[40 * mm, 30 * mm])
        t.setStyle(TableStyle([("FONTNAME", (0, 0), (-1, -1), font), ("GRID", (0, 0), (-1, -1), 0.5, colors.grey)]))
        story.append(t)

    save_flow("PDF20_table_only.pdf", pdf20)

    return names


def main():
    assets = make_assets()
    o = gen_office_wave4(assets)
    p = gen_pdf(assets)
    print(f"office_w4={len(o)} pdf={len(p)} total={len(o)+len(p)}")


if __name__ == "__main__":
    main()
